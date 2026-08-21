from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Colors (from real CSS) ──
BG    = RGBColor(0x0a, 0x0e, 0x1a)
BG2   = RGBColor(0x0f, 0x15, 0x26)
SURF  = RGBColor(0x14, 0x1b, 0x30)
SURF2 = RGBColor(0x1b, 0x24, 0x40)
ACC   = RGBColor(0x4f, 0x7d, 0xff)
ACC2  = RGBColor(0x7a, 0x9d, 0xff)
GRN   = RGBColor(0x2f, 0xd0, 0x7f)
RED   = RGBColor(0xff, 0x5c, 0x6c)
AMB   = RGBColor(0xff, 0xb4, 0x54)
ORG   = RGBColor(0xff, 0x8a, 0x4c)
PUR   = RGBColor(0xb0, 0x7a, 0xff)
TXT   = RGBColor(0xee, 0xf2, 0xff)
TXT2  = RGBColor(0xd4, 0xdb, 0xf2)
MUT   = RGBColor(0x8b, 0x95, 0xb8)
WHT   = RGBColor(0xff, 0xff, 0xff)
GGRAY = RGBColor(0x56, 0x5f, 0x85)
RED_BG  = RGBColor(0x3a, 0x15, 0x1c)
BLUE_BG = RGBColor(0x14, 0x20, 0x4a)
GRN_BG  = RGBColor(0x0f, 0x2a, 0x1c)
AMB_BG  = RGBColor(0x3a, 0x2e, 0x15)

# ── Layout ──
PX, PY, PW, PH = 0.3, 0.7, 4.0, 6.8
R1X, R1W = 4.8, 4.0
R2X, R2W = 9.1, 3.7

# ── Helpers ──
def _bg(sl, c):
    sl.background.fill.solid(); sl.background.fill.fore_color.rgb = c

def _bar(sl):
    for y in [0, 7.42]:
        b = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(y), Inches(13.333), Inches(0.07))
        b.fill.solid(); b.fill.fore_color.rgb = ACC; b.line.fill.background()

def _slide(sl):
    _bg(sl, BG); _bar(sl)

def _title(sl, t, sub=None):
    box = sl.shapes.add_textbox(Inches(0.7), Inches(0.15), Inches(12), Inches(0.55))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t; p.font.size = Pt(28); p.font.color.rgb = TXT; p.font.bold = True
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub; p2.font.size = Pt(14); p2.font.color.rgb = MUT; p2.space_before = Pt(4)

def _txt(sl, l, t, w, h, text, sz=14, c=TXT, b=False, al=PP_ALIGN.LEFT):
    box = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(sz); p.font.color.rgb = c; p.font.bold = b; p.alignment = al
    return tf

def _card(sl, l, t, w, h, fill=SURF):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    return sh

def _circ(sl, l, t, d, fill):
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    return sh

def _rect(sl, l, t, w, h, fill):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill; sh.line.fill.background()
    return sh

def _pill(sl, l, t, text, tc, bc, w=None, sz=9):
    w = w or max(0.7, len(text) * 0.1 + 0.22)
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.22))
    sh.fill.solid(); sh.fill.fore_color.rgb = bc; sh.line.fill.background()
    _txt(sl, l, t+0.02, w, 0.18, text, sz, tc, True, PP_ALIGN.CENTER)

def _mbadge(sl, l, t, text, tc, bc):
    w = max(0.6, len(text)*0.08+0.15)
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.2))
    sh.fill.solid(); sh.fill.fore_color.rgb = bc; sh.line.fill.background()
    _txt(sl, l, t+0.01, w, 0.18, text, 8, tc, True, PP_ALIGN.CENTER)
    return w

def _rline(sl, y, col):
    _rect(sl, R1X, y, 0.06, 0.35, col)

# ── Phone: returns absolute screen coords ──
def _phone(sl):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(PX+0.03), Inches(PY+0.03), Inches(PW), Inches(PH))
    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(0x05, 0x08, 0x12); sh.line.fill.background()
    body = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(PX), Inches(PY), Inches(PW), Inches(PH))
    body.fill.solid(); body.fill.fore_color.rgb = RGBColor(0x0f, 0x15, 0x26); body.line.color.rgb = RGBColor(0x2a, 0x3a, 0x5a); body.line.width = Pt(0.5)
    sx = PX + 0.07; sy = PY + 0.07; sw = PW - 0.14; sh_h = PH - 0.14
    scr = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx), Inches(sy), Inches(sw), Inches(sh_h))
    scr.fill.solid(); scr.fill.fore_color.rgb = BG; scr.line.fill.background()
    nw = 0.9; nx = PX + (PW - nw) / 2
    notch = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(nx), Inches(PY+0.01), Inches(nw), Inches(0.04))
    notch.fill.solid(); notch.fill.fore_color.rgb = RGBColor(0x0f, 0x15, 0x26); notch.line.fill.background()
    return (sx, sy, sw, sh_h)

def _topbar(sl, sx, sy, sw):
    _rect(sl, sx, sy, sw, 0.42, RGBColor(0x0e, 0x12, 0x22))
    _circ(sl, sx+0.1, sy+0.08, 0.26, ACC)
    _txt(sl, sx+0.1, sy+0.09, 0.26, 0.18, "N", 10, WHT, True, PP_ALIGN.CENTER)
    _txt(sl, sx+0.42, sy+0.08, 1.0, 0.22, "NexAlert", 11, TXT, True)
    _circ(sl, sx+sw-0.7, sy+0.08, 0.26, ACC)
    _txt(sl, sx+sw-0.7, sy+0.08, 0.26, 0.22, "+", 12, WHT, True, PP_ALIGN.CENTER)
    _txt(sl, sx+sw-0.34, sy+0.08, 0.28, 0.22, "...", 11, MUT, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 1 - PORTADA
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
ov = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
ov.fill.solid(); ov.fill.fore_color.rgb = RGBColor(0x06, 0x0a, 0x14); ov.line.fill.background()
_circ(sl, 6.15, 0.55, 1.1, ACC)
_txt(sl, 6.15, 0.67, 1.1, 0.7, "N", 34, WHT, True, PP_ALIGN.CENTER)
_txt(sl, 1, 1.95, 11.3, 0.8, "CAPACITACION NEXALERT", 44, WHT, True, PP_ALIGN.CENTER)
_txt(sl, 1, 2.95, 11.3, 0.5, "Guia completa para tecnicos en campo", 20, MUT, False, PP_ALIGN.CENTER)
d = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(3.65), Inches(2.7), Inches(0.03))
d.fill.solid(); d.fill.fore_color.rgb = ACC; d.line.fill.background()
_txt(sl, 1, 3.9, 11.3, 0.4, "App Android  |  v1.5.5", 16, ACC2, False, PP_ALIGN.CENTER)
feats = ["Crear reportes", "Tomar fotos", "Firma digital", "Trabajar offline", "GPS automatico", "Compartir"]
for i, f in enumerate(feats):
    _pill(sl, 2.1 + i*1.65, 4.7, f, ACC2, SURF2, 1.5, 10)
_txt(sl, 1, 5.7, 11.3, 0.4, "NEXUS  |  2026", 13, MUT, False, PP_ALIGN.CENTER)
print("Slide 1 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 2 - LOGIN
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Bienvenido - Primeros Pasos", "Configura tu acceso y comienza a reportar")
sx, sy, sw, sh = _phone(sl)

# Login content inside phone (absolute coords from sy)
_circ(sl, sx+sw/2-0.35, sy+0.55, 0.7, ACC)
_txt(sl, sx+sw/2-0.35, sy+0.63, 0.7, 0.5, "N", 20, WHT, True, PP_ALIGN.CENTER)
_txt(sl, sx+0.2, sy+1.35, sw-0.4, 0.25, "Bienvenido, tecnico", 11, TXT, True, PP_ALIGN.CENTER)
# Inputs
_card(sl, sx+0.2, sy+1.8, sw-0.4, 0.4, BG2)
_txt(sl, sx+0.3, sy+1.83, sw-0.6, 0.16, "Usuario", 8, MUT)
_card(sl, sx+0.2, sy+2.35, sw-0.4, 0.4, BG2)
_txt(sl, sx+0.3, sy+2.38, sw-0.6, 0.16, "Contrasena", 8, MUT)
# Buttons
sh2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.2), Inches(sy+2.9), Inches(sw-0.4), Inches(0.38))
sh2.fill.solid(); sh2.fill.fore_color.rgb = ACC; sh2.line.fill.background()
_txt(sl, sx+0.2, sy+2.92, sw-0.4, 0.34, "Entrar", 11, WHT, True, PP_ALIGN.CENTER)
sh3 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.2), Inches(sy+3.4), Inches(sw-0.4), Inches(0.35))
sh3.fill.solid(); sh3.fill.fore_color.rgb = SURF2; sh3.line.color.rgb = RGBColor(0x2a, 0x3a, 0x5a); sh3.line.width = Pt(0.5)
_txt(sl, sx+0.2, sy+3.43, sw-0.4, 0.3, "Entrar con biometria", 10, TXT, True, PP_ALIGN.CENTER)

# Config card
_card(sl, R1X, 0.85, R1W, 5.9)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Configuracion inicial", 16, ACC, True)
steps = [
    ("1", "Descargar la app", "Instala NexAlert desde el enlace\nde tu administrador", ACC),
    ("2", "Iniciar sesion", "Usa tu usuario y contrasena\nasignados por el gerente", GRN),
    ("3", "Configurar biometria", "Activa huella o rostro para\nacceso rapido futuro", PUR),
    ("4", "Primer reporte", "Toca el boton + y crea tu\nprimer reporte de falla", AMB),
]
for i, (num, title_t, desc, col) in enumerate(steps):
    y = 1.45 + i*1.3
    _circ(sl, R1X+0.2, y+0.05, 0.4, col)
    _txt(sl, R1X+0.2, y+0.08, 0.4, 0.3, num, 14, WHT, True, PP_ALIGN.CENTER)
    _txt(sl, R1X+0.75, y+0.02, 3.0, 0.22, title_t, 13, TXT, True)
    _txt(sl, R1X+0.75, y+0.28, 3.0, 0.5, desc, 11, MUT)

# Security card
_card(sl, R2X, 0.85, R2W, 5.9)
_txt(sl, R2X+0.2, 0.95, R2W-0.4, 0.3, "Seguridad", 16, RED, True)
sec_items = [
    ("Cada tecnico tiene usuario propio", ACC),
    ("La biometria es opcional pero recomendada", GRN),
    ("Las credenciales se guardan localmente", PUR),
    ("Puedes cerrar sesion desde el menu", AMB),
    ("La sesion se cierra al cerrar la app", RED),
    ("Los reportes incluyen tu identificacion", ACC),
    ("Tu ubicacion GPS se registra al reportar", GRN),
]
for i, (text, col) in enumerate(sec_items):
    y = 1.5 + i*0.7
    _circ(sl, R2X+0.2, y+0.1, 0.18, col)
    _txt(sl, R2X+0.5, y+0.05, R2W-0.8, 0.45, text, 11, TXT2)
print("Slide 2 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 3 - PANTALLA PRINCIPAL
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Pantalla Principal", "Tu lista de reportes asignados con filtros y busqueda")
sx, sy, sw, sh = _phone(sl)
_topbar(sl, sx, sy, sw)

# Content starts at sy+0.50 (below topbar)
_txt(sl, sx+0.15, sy+0.52, sw-0.3, 0.25, "Mis reportes", 13, TXT, True)
_txt(sl, sx+0.15, sy+0.77, sw-0.3, 0.15, "Actualizado hace 2 min", 8, MUT)
# Search
_card(sl, sx+0.12, sy+1.0, sw-0.24, 0.32, BG2)
_txt(sl, sx+0.22, sy+1.03, sw-0.44, 0.22, "Buscar cliente, equipo...", 9, MUT)
# Filters row
_card(sl, sx+0.12, sy+1.4, 1.15, 0.28, BG2)
_txt(sl, sx+0.2, sy+1.42, 0.95, 0.2, "Todos", 9, MUT)
_card(sl, sx+1.35, sy+1.4, sw-1.47, 0.28, BG2)
_txt(sl, sx+1.43, sy+1.42, sw-1.6, 0.2, "Prioridad", 9, MUT)

# Report card 1
y1 = sy + 1.85
c1 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(y1), Inches(sw-0.24), Inches(1.05))
c1.fill.solid(); c1.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x30); c1.line.color.rgb = RGBColor(0x2a, 0x3a, 0x5a); c1.line.width = Pt(0.3)
_rect(sl, sx+0.12, y1, 0.03, 1.05, RED)
_txt(sl, sx+0.24, y1+0.06, 1.8, 0.2, "Minera del Sur", 10, TXT, True)
_mbadge(sl, sx+sw-0.82, y1+0.06, "Abierto", WHT, RED_BG)
_txt(sl, sx+0.24, y1+0.3, sw-0.44, 0.15, "Compresor Atlas Copco", 8, MUT)
_txt(sl, sx+0.24, y1+0.5, sw-0.44, 0.3, "Fuga de aceite en el compresor principal, requiere revision urgente", 8, TXT2)
_mbadge(sl, sx+0.24, y1+0.82, "Urgente", WHT, RED)
_mbadge(sl, sx+0.82, y1+0.82, "Fotos", WHT, BLUE_BG)
_txt(sl, sx+sw-0.8, y1+0.84, 0.65, 0.15, "15 Ago", 7, MUT, False, PP_ALIGN.RIGHT)

# Report card 2
y2 = sy + 3.1
c2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(y2), Inches(sw-0.24), Inches(1.05))
c2.fill.solid(); c2.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x30); c2.line.color.rgb = RGBColor(0x2a, 0x3a, 0x5a); c2.line.width = Pt(0.3)
_rect(sl, sx+0.12, y2, 0.03, 1.05, ACC)
_txt(sl, sx+0.24, y2+0.06, 1.8, 0.2, "Cementos Argos", 10, TXT, True)
_mbadge(sl, sx+sw-0.9, y2+0.06, "En proceso", WHT, BLUE_BG)
_txt(sl, sx+0.24, y2+0.3, sw-0.44, 0.15, "Molino vertical VRM-4", 8, MUT)
_txt(sl, sx+0.24, y2+0.5, sw-0.44, 0.3, "Vibracion excesiva en rodamiento superior, requiere cambio", 8, TXT2)
_mbadge(sl, sx+0.24, y2+0.82, "Normal", WHT, ACC)
_txt(sl, sx+sw-0.8, y2+0.84, 0.65, 0.15, "14 Ago", 7, MUT, False, PP_ALIGN.RIGHT)

# Report card 3
y3 = sy + 4.35
c3 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(y3), Inches(sw-0.24), Inches(1.05))
c3.fill.solid(); c3.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x30); c3.line.color.rgb = RGBColor(0x2a, 0x3a, 0x5a); c3.line.width = Pt(0.3)
_rect(sl, sx+0.12, y3, 0.03, 1.05, GRN)
_txt(sl, sx+0.24, y3+0.06, 1.8, 0.2, "Ecopetrol SA", 10, TXT, True)
_mbadge(sl, sx+sw-0.75, y3+0.06, "Resuelto", WHT, GRN_BG)
_txt(sl, sx+0.24, y3+0.3, sw-0.44, 0.15, "Turbina de gas TG-200", 8, MUT)
_txt(sl, sx+0.24, y3+0.5, sw-0.44, 0.3, "Fallo en sistema de lubricacion, repuesto instalado", 8, TXT2)
_mbadge(sl, sx+0.24, y3+0.82, "Baja", WHT, GGRAY)
_txt(sl, sx+sw-0.8, y3+0.84, 0.65, 0.15, "10 Ago", 7, MUT, False, PP_ALIGN.RIGHT)

# Right - Elements
_card(sl, R1X, 0.85, R1W, 5.9)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Elementos de la pantalla", 16, ACC, True)
elems = [
    ("Topbar", "Logo + boton nuevo (+) + menu (...)", ACC),
    ("Busqueda", "Escribe nombre del cliente, equipo o problema", GRN),
    ("Filtro estado", "Abierto, En proceso, Resuelto, En espera", PUR),
    ("Filtro prioridad", "Urgente, Alta, Normal, Baja", RED),
    ("Tarjetas", "Cada reporte se muestra como tarjeta", ACC),
    ("Borde de color", "Indica el estado del reporte", GRN),
    ("Badges", "Estado y prioridad como etiquetas", PUR),
    ("Fecha", "Dia de creacion del reporte", MUT),
]
for i, (label, desc, col) in enumerate(elems):
    y = 1.4 + i*0.55
    _rect(sl, R1X+0.2, y+0.05, 0.06, 0.38, col)
    _txt(sl, R1X+0.4, y+0.02, 1.4, 0.2, label, 11, col, True)
    _txt(sl, R1X+1.9, y+0.02, 1.9, 0.38, desc, 10, MUT)

# Right - Legend
_card(sl, R2X, 0.85, R2W, 5.9)
_txt(sl, R2X+0.2, 0.95, R2W-0.4, 0.3, "Leyenda de colores", 16, ACC, True)
legend = [
    ("Abierto", "Reporte nuevo sin atender", WHT, RED_BG),
    ("En proceso", "Siendo atendido por tecnico", WHT, BLUE_BG),
    ("Espera", "Espera repuesto o cliente", RGBColor(0x1a, 0x12, 0x00), AMB_BG),
    ("Resuelto", "Problema solucionado", WHT, GRN_BG),
    ("Urgente", "Accion inmediata requerida", WHT, RED),
    ("Alta", "Alta prioridad", WHT, ORG),
    ("Normal", "Prioridad estandar", WHT, ACC),
    ("Baja", "Mantenimiento preventivo", WHT, GGRAY),
]
for i, (label, desc, tc, bc) in enumerate(legend):
    y = 1.4 + i*0.58
    _pill(sl, R2X+0.2, y+0.02, label, tc, bc, 0.85, 9)
    _txt(sl, R2X+1.15, y+0.02, 2.3, 0.4, desc, 10, MUT)
print("Slide 3 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 4 - FILTRAR
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Filtrar y Buscar Reportes", "Encuentra rapidamente el reporte que necesitas")
sx, sy, sw, sh = _phone(sl)
_topbar(sl, sx, sy, sw)
_txt(sl, sx+0.15, sy+0.52, sw-0.3, 0.25, "Mis reportes", 13, TXT, True)
# Search
_card(sl, sx+0.12, sy+0.87, sw-0.24, 0.32, BG2)
_txt(sl, sx+0.22, sy+0.9, sw-0.44, 0.22, "Compresor", 9, TXT)
# Estado dropdown expanded
_card(sl, sx+0.12, sy+1.27, 1.35, 0.28, BG2)
_txt(sl, sx+0.2, sy+1.29, 1.15, 0.2, "Abierto", 9, RED)
_card(sl, sx+0.12, sy+1.57, 1.35, 1.45, SURF2)
for i, (est, col) in enumerate([("Todos", MUT), ("Abierto", RED), ("En proceso", ACC), ("En espera", AMB), ("Resuelto", GRN)]):
    _txt(sl, sx+0.22, sy+1.62+i*0.26, 1.1, 0.22, est, 9, col)
# Prioridad
_card(sl, sx+1.55, sy+1.27, sw-1.67, 0.28, BG2)
_txt(sl, sx+1.63, sy+1.29, sw-1.8, 0.2, "Prioridad", 9, MUT)
# Date filters (below dropdown)
_card(sl, sx+0.12, sy+3.12, 1.25, 0.24, BG2)
_txt(sl, sx+0.2, sy+3.14, 1.1, 0.18, "Desde: 01/08/2026", 7, MUT)
_card(sl, sx+1.45, sy+3.12, sw-1.57, 0.24, BG2)
_txt(sl, sx+1.53, sy+3.14, sw-1.7, 0.18, "Hasta: 20/08/2026", 7, MUT)
# Filtered result
y4 = sy + 3.5
c4 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(y4), Inches(sw-0.24), Inches(0.9))
c4.fill.solid(); c4.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x30); c4.line.color.rgb = RGBColor(0x2a, 0x3a, 0x5a); c4.line.width = Pt(0.3)
_rect(sl, sx+0.12, y4, 0.03, 0.9, RED)
_txt(sl, sx+0.24, y4+0.06, 1.8, 0.2, "Minera del Sur", 10, TXT, True)
_mbadge(sl, sx+sw-0.82, y4+0.06, "Abierto", WHT, RED_BG)
_txt(sl, sx+0.24, y4+0.28, sw-0.44, 0.15, "Compresor Atlas Copco XAS 185", 8, MUT)
_txt(sl, sx+0.24, y4+0.46, sw-0.44, 0.32, "Fuga de aceite en el compresor principal, requiere revision urgente del sello", 8, TXT2)
_txt(sl, sx+0.24, y4+0.7, 2.0, 0.15, "1 resultado encontrado", 8, GRN, True)

# Right
_card(sl, R1X, 0.85, R1W, 5.9)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Como filtrar", 16, ACC, True)
filts = [
    ("Busqueda por texto", "Escribe en el campo de busqueda. Busca por nombre de cliente, nombre del equipo o descripcion del problema. Se filtra en tiempo real.", ACC),
    ("Filtro por estado", "Selecciona un estado del dropdown: Todos, Abierto, En proceso, En espera, o Resuelto.", GRN),
    ("Filtro por prioridad", "Filtra por nivel de urgencia: Urgente, Alta, Normal o Baja.", RED),
    ("Filtros de fecha", "Selecciona un rango Desde/Hasta para ver reportes de un periodo especifico.", AMB),
    ("Combinar filtros", "Puedes usar varios filtros al mismo tiempo.", PUR),
    ("Limpiar filtros", "Selecciona 'Todos' en cada dropdown para quitar los filtros.", MUT),
]
for i, (title_t, desc, col) in enumerate(filts):
    y = 1.4 + i*0.85
    _rect(sl, R1X+0.2, y+0.05, 0.06, 0.7, col)
    _txt(sl, R1X+0.4, y+0.02, R1W-0.7, 0.2, title_t, 12, col, True)
    _txt(sl, R1X+0.4, y+0.25, R1W-0.7, 0.5, desc, 10, MUT)
print("Slide 4 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 5 - CREAR REPORTE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Crear un Nuevo Reporte", "Toca el boton + y completa el formulario")
sx, sy, sw, sh = _phone(sl)
_topbar(sl, sx, sy, sw)
_txt(sl, sx+0.15, sy+0.52, 1.2, 0.22, "< Volver", 9, ACC, True)
# Form card
fc = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(sy+0.82), Inches(sw-0.24), Inches(5.5))
fc.fill.solid(); fc.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x30); fc.line.color.rgb = RGBColor(0x2a, 0x3a, 0x5a); fc.line.width = Pt(0.3)
_txt(sl, sx+0.24, sy+0.92, sw-0.48, 0.24, "Nuevo reporte", 12, TXT, True)
fields = [("Cliente *", "Minera del Sur"), ("Equipo", "Compresor Atlas Copco"), ("Descripcion *", "Fuga de aceite en el\ncompresor principal...")]
for i, (label, val) in enumerate(fields):
    fy = sy + 1.25 + i*0.9
    _txt(sl, sx+0.24, fy, sw-0.48, 0.16, label, 9, MUT)
    _card(sl, sx+0.24, fy+0.18, sw-0.48, 0.5, BG2)
    _txt(sl, sx+0.34, fy+0.22, sw-0.68, 0.42, val, 10, TXT2)
_txt(sl, sx+0.24, sy+3.98, sw-0.48, 0.16, "Prioridad", 9, MUT)
_card(sl, sx+0.24, sy+4.18, sw-0.48, 0.35, BG2)
_txt(sl, sx+0.34, sy+4.21, sw-0.68, 0.28, "Normal", 10, ACC)
sh2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.24), Inches(sy+4.7), Inches(sw-0.48), Inches(0.38))
sh2.fill.solid(); sh2.fill.fore_color.rgb = ACC; sh2.line.fill.background()
_txt(sl, sx+0.24, sy+4.72, sw-0.48, 0.34, "Guardar reporte", 11, WHT, True, PP_ALIGN.CENTER)
_txt(sl, sx+0.24, sy+5.2, sw-0.48, 0.3, "GPS capturado automaticamente", 8, GRN, True, PP_ALIGN.CENTER)

# Steps
_card(sl, R1X, 0.85, R1W, 2.8)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Pasos para crear", 16, ACC, True)
for i, (step, desc, col) in enumerate([
    ("1. Toca el boton +", "Se abre el formulario nuevo reporte", ACC),
    ("2. Escribe el cliente", "Nombre de la empresa que reporta", GRN),
    ("3. Describe el problema", "Se claro y especifico sobre la falla", PUR),
    ("4. Selecciona prioridad", "Normal por defecto, cambia si es urgente", AMB),
    ("5. Toca Guardar", "El reporte se envia al servidor", GRN),
]):
    y = 1.4 + i*0.44
    _circ(sl, R1X+0.2, y+0.03, 0.3, col)
    _txt(sl, R1X+0.2, y+0.05, 0.3, 0.2, str(i+1), 9, WHT, True, PP_ALIGN.CENTER)
    _txt(sl, R1X+0.6, y, 3.2, 0.17, step, 11, TXT, True)
    _txt(sl, R1X+0.6, y+0.18, 3.2, 0.22, desc, 10, MUT)

_card(sl, R1X, 3.85, R1W, 2.8)
_txt(sl, R1X+0.2, 3.95, R1W-0.4, 0.3, "Campos obligatorios", 14, RED, True)
for i, (campo, desc) in enumerate([
    ("Cliente", "Nombre de la empresa. Escribe y el sistema autocompleta."),
    ("Descripcion", "Detalla el problema lo mas claro posible."),
]):
    y = 4.4 + i*0.9
    _txt(sl, R1X+0.2, y, R1W-0.4, 0.2, campo, 12, RED, True)
    _txt(sl, R1X+0.2, y+0.28, R1W-0.4, 0.5, desc, 10, MUT)

_card(sl, R2X, 0.85, R2W, 5.9)
_txt(sl, R2X+0.2, 0.95, R2W-0.4, 0.3, "Consejos", 16, GRN, True)
tips = [
    ("GPS automatico", "La ubicacion se captura solo. No necesitas hacer nada.", GRN),
    ("Descripcion clara", "Incluye marca, modelo y sintomas exactos.", ACC),
    ("Prioridad real", "Solo marca urgente si hay parada total.", RED),
    ("Sin fotos aun", "Puedes agregar fotos despues desde el detalle.", MUT),
    ("Modo offline", "Si no tienes internet, se guarda local y se envia despues.", PUR),
]
for i, (tip, desc, col) in enumerate(tips):
    y = 1.4 + i*0.95
    _rect(sl, R2X+0.2, y+0.05, 0.06, 0.8, col)
    _txt(sl, R2X+0.4, y+0.02, R2W-0.7, 0.2, tip, 11, col, True)
    _txt(sl, R2X+0.4, y+0.28, R2W-0.7, 0.5, desc, 10, MUT)
print("Slide 5 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 6 - DETALLE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Detalle del Reporte", "Toca una tarjeta para ver toda la informacion")
sx, sy, sw, sh = _phone(sl)
_topbar(sl, sx, sy, sw)
_txt(sl, sx+0.15, sy+0.52, 1.2, 0.22, "< Volver", 9, ACC, True)
# Header
dh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(sy+0.82), Inches(sw-0.24), Inches(1.3))
dh.fill.solid(); dh.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x30); dh.line.color.rgb = RGBColor(0x2a, 0x3a, 0x5a); dh.line.width = Pt(0.3)
_txt(sl, sx+0.24, sy+0.9, 2.0, 0.22, "Minera del Sur", 12, TXT, True)
_mbadge(sl, sx+sw-0.85, sy+0.9, "Abierto", WHT, RED_BG)
_txt(sl, sx+0.24, sy+1.16, 2.2, 0.16, "Compresor Atlas Copco XAS 185", 8, MUT)
_mbadge(sl, sx+0.24, sy+1.36, "Urgente", WHT, RED)
_mbadge(sl, sx+0.84, sy+1.36, "GPS", WHT, GRN_BG)
_txt(sl, sx+0.24, sy+1.62, sw-0.48, 0.14, "PROBLEMA", 7, MUT, True)
_txt(sl, sx+0.24, sy+1.78, sw-0.48, 0.28, "Fuga de aceite en compresor principal, sello danado", 8, TXT2)
# Fields
for i, (lbl, val, col) in enumerate([("ESTADO", "En proceso", ACC), ("PRIORIDAD", "Urgente", RED), ("TECNICO", "Carlos Rodriguez", MUT), ("CREADO", "15/08/2026 09:30", MUT)]):
    fy = sy + 2.2 + i*0.28
    _txt(sl, sx+0.24, fy, 1.1, 0.14, lbl, 6, MUT, True)
    _txt(sl, sx+1.35, fy, 1.6, 0.14, val, 8, col)
# State buttons
_txt(sl, sx+0.15, sy+3.35, sw-0.3, 0.14, "Cambiar estado:", 7, MUT, True)
estados = [("Abierto", RED), ("En pro.", ACC), ("Resuelto", GRN), ("Repuesto", AMB), ("Cliente", AMB)]
for i, (est, col) in enumerate(estados):
    ex = sx + 0.12 + i*(sw-0.24)/5
    ew = (sw-0.24)/5 - 0.02
    is_active = (i == 1)
    ec = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(ex), Inches(sy+3.52), Inches(ew), Inches(0.26))
    ec.fill.solid(); ec.fill.fore_color.rgb = ACC if is_active else SURF2; ec.line.fill.background()
    _txt(sl, ex, sy+3.54, ew, 0.22, est, 6, WHT if is_active else MUT, is_active, PP_ALIGN.CENTER)
# Photos
_txt(sl, sx+0.15, sy+3.88, sw-0.3, 0.14, "FOTOGRAFIAS", 7, MUT, True)
for i in range(3):
    fx = sx + 0.12 + i*(sw-0.24)/3
    fw = (sw-0.24)/3 - 0.04
    fc2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(fx), Inches(sy+4.06), Inches(fw), Inches(0.55))
    fc2.fill.solid(); fc2.fill.fore_color.rgb = SURF2; fc2.line.fill.background()
_txt(sl, sx+0.12, sy+4.7, (sw-0.24)/2-0.02, 0.28, "Tomar foto", 9, WHT, True, PP_ALIGN.CENTER)
_txt(sl, sx+0.12+(sw-0.24)/2+0.02, sy+4.7, (sw-0.24)/2-0.02, 0.28, "Galeria", 9, TXT, True, PP_ALIGN.CENTER)
# Firma
_txt(sl, sx+0.15, sy+5.08, sw-0.3, 0.14, "FIRMA DIGITAL", 7, MUT, True)
_firma = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(sy+5.26), Inches(sw-0.24), Inches(0.5))
_firma.fill.solid(); _firma.fill.fore_color.rgb = BG2; _firma.line.fill.background()
_txt(sl, sx+0.12, sy+5.4, sw-0.24, 0.22, "Dibuja tu firma aqui", 8, MUT, False, PP_ALIGN.CENTER)
# Notas
_txt(sl, sx+0.15, sy+5.86, sw-0.3, 0.14, "NOTAS", 7, MUT, True)
_card(sl, sx+0.12, sy+6.04, sw-0.24, 0.4, BG2)
_txt(sl, sx+0.22, sy+6.08, sw-0.48, 0.3, "Sello danado, requiere repuesto...", 7, TXT2)

# Right
_card(sl, R1X, 0.85, R1W, 5.9)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Secciones del detalle", 16, ACC, True)
sections = [
    ("Cabecera", "Cliente, equipo, badges de estado y prioridad", ACC),
    ("Problema", "Descripcion detallada del fallo reportado", RED),
    ("Estado", "Tu estado actual y botones para cambiarlo", AMB),
    ("Ubicacion GPS", "Coordenadas del lugar del fallo", GRN),
    ("Fotografias", "Fotos tomadas en campo (camara o galeria)", ACC),
    ("Firma digital", "Firma del cliente o confirmacion de entrega", PUR),
    ("Notas internas", "Comentarios privados sobre el reporte", MUT),
    ("Compartir", "Enviar resumen por WhatsApp o email", GRN),
]
for i, (sec, desc, col) in enumerate(sections):
    y = 1.4 + i*0.55
    _rect(sl, R1X+0.2, y+0.05, 0.06, 0.4, col)
    _txt(sl, R1X+0.4, y+0.02, 1.5, 0.2, sec, 11, col, True)
    _txt(sl, R1X+2.0, y+0.02, 1.8, 0.4, desc, 10, MUT)
print("Slide 6 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 7 - ESTADO
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Cambiar Estado del Reporte", "Actualiza el estado a medida que avanzas en la reparacion")
_card(sl, 0.5, 0.9, 12.3, 1.8)
_txt(sl, 0.7, 1.0, 3.0, 0.3, "Flujo de estados", 16, ACC, True)
for i, (est, desc, col) in enumerate([
    ("abierto", "Reporte nuevo,\nesperando atencion", RED),
    ("en_proceso", "Trabajando\nen la falla", ACC),
    ("resuelto", "Problema\nsolucionado", GRN),
]):
    x = 1.0 + i*4.2
    sh2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.45), Inches(2.8), Inches(1.05))
    sh2.fill.solid(); sh2.fill.fore_color.rgb = SURF2; sh2.line.color.rgb = col; sh2.line.width = Pt(1.5)
    _txt(sl, x+0.1, 1.5, 2.6, 0.28, est, 14, col, True, PP_ALIGN.CENTER)
    _txt(sl, x+0.1, 1.82, 2.6, 0.55, desc, 11, MUT, False, PP_ALIGN.CENTER)
    if i < 2:
        arr = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x+2.95), Inches(1.8), Inches(1.1), Inches(0.22))
        arr.fill.solid(); arr.fill.fore_color.rgb = MUT; arr.line.fill.background()

_card(sl, 0.5, 3.0, 6.1, 3.8)
_txt(sl, 0.7, 3.1, 5.7, 0.3, "Estados de espera", 15, AMB, True)
for i, (est, desc, col) in enumerate([
    ("espera_repuesto", "Necesitas una pieza para continuar. Describe que repuesto falta en las notas.", AMB),
    ("espera_cliente", "Requieres informacion adicional del cliente para avanzar.", AMB),
]):
    y = 3.6 + i*1.3
    _card(sl, 0.7, y, 5.7, 1.1, SURF2)
    _rect(sl, 0.7, y, 0.06, 1.1, col)
    _pill(sl, 0.95, y+0.12, est, TXT, col, 2.0, 9)
    _txt(sl, 0.95, y+0.45, 5.2, 0.55, desc, 11, MUT)

_card(sl, 6.8, 3.0, 6.0, 3.8)
_txt(sl, 7.0, 3.1, 5.6, 0.3, "Prioridades", 15, RED, True)
for i, (pri, desc, col) in enumerate([
    ("baja", "Mantenimiento preventivo. Se agenda para revision periodica.", GGRAY),
    ("normal", "Fallo estandar que afecta operacion. Se atiende en horas normales.", ACC),
    ("alta", "Fallo critico que requiere atencion rapida. Se prioriza.", ORG),
    ("urgente", "Parada total de maquinaria. Atencion inmediata, notifica a todos.", RED),
]):
    y = 3.6 + i*0.82
    _card(sl, 7.0, y, 5.6, 0.72, SURF2)
    _rect(sl, 7.0, y, 0.06, 0.72, col)
    _pill(sl, 7.2, y+0.1, pri, TXT if col==GGRAY else WHT, col, 0.95, 9)
    _txt(sl, 8.3, y+0.08, 4.1, 0.55, desc, 10, MUT)
print("Slide 7 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 8 - FOTOS
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Fotografias del Reporte", "Documenta visualmente el problema con fotos")
sx, sy, sw, sh = _phone(sl)
_topbar(sl, sx, sy, sw)
_txt(sl, sx+0.15, sy+0.52, 1.2, 0.22, "< Volver", 9, ACC, True)
_txt(sl, sx+0.15, sy+0.82, sw-0.3, 0.22, "FOTOGRAFIAS", 9, MUT, True)
for r in range(2):
    for c in range(3):
        fx = sx + 0.12 + c*(sw-0.24)/3
        fy = sy + 1.12 + r*1.0
        fw = (sw-0.24)/3 - 0.04
        fc2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(fx), Inches(fy), Inches(fw), Inches(0.88))
        fc2.fill.solid(); fc2.fill.fore_color.rgb = SURF2; fc2.line.fill.background()
        _txt(sl, fx, fy+0.35, fw, 0.18, "foto " + str(r*3+c+1), 8, MUT, False, PP_ALIGN.CENTER)
_txt(sl, sx+0.12, sy+3.18, (sw-0.24)/2-0.02, 0.3, "Tomar foto", 10, WHT, True, PP_ALIGN.CENTER)
_txt(sl, sx+0.12+(sw-0.24)/2+0.02, sy+3.18, (sw-0.24)/2-0.02, 0.3, "Galeria", 10, TXT, True, PP_ALIGN.CENTER)

_card(sl, R1X, 0.85, R1W, 5.9)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Como tomar fotos", 16, ACC, True)
for i, (title_t, desc, col) in enumerate([
    ("Camara", "Toca 'Tomar foto' para abrir la camara del dispositivo.", ACC),
    ("Galeria", "Toca 'Galeria' para seleccionar una foto existente.", GRN),
    ("Multiples fotos", "Puedes agregar varias fotos por reporte.", PUR),
    ("Eliminar foto", "Toca una foto para verla completa. Usa la basura para eliminarla.", RED),
    ("Consejos", "Foto del dano, del equipo completo, del numero de serie.", AMB),
    ("Tamano", "Las fotos se comprimen automaticamente antes de enviarse.", MUT),
]):
    y = 1.4 + i*0.85
    _rect(sl, R1X+0.2, y+0.05, 0.06, 0.7, col)
    _txt(sl, R1X+0.4, y+0.02, R1W-0.7, 0.2, title_t, 12, col, True)
    _txt(sl, R1X+0.4, y+0.28, R1W-0.7, 0.5, desc, 10, MUT)
print("Slide 8 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 9 - FIRMA
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Firma Digital", "Registra la confirmacion de entrega con tu firma")
sx, sy, sw, sh = _phone(sl)
_topbar(sl, sx, sy, sw)
_txt(sl, sx+0.15, sy+0.52, 1.2, 0.22, "< Volver", 9, ACC, True)
_txt(sl, sx+0.15, sy+0.82, sw-0.3, 0.22, "FIRMA DIGITAL", 9, MUT, True)
fc = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(sy+1.12), Inches(sw-0.24), Inches(2.2))
fc.fill.solid(); fc.fill.fore_color.rgb = BG2; fc.line.fill.background()
_txt(sl, sx+0.12, sy+2.0, sw-0.24, 0.22, "Dibuja tu firma aqui", 10, MUT, False, PP_ALIGN.CENTER)
_txt(sl, sx+0.12, sy+3.5, (sw-0.24)/2-0.02, 0.3, "Limpiar", 10, MUT, True, PP_ALIGN.CENTER)
_txt(sl, sx+0.12+(sw-0.24)/2+0.02, sy+3.5, (sw-0.24)/2-0.02, 0.3, "Guardar firma", 10, WHT, True, PP_ALIGN.CENTER)

_card(sl, R1X, 0.85, R1W, 5.9)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Sobre la firma digital", 16, ACC, True)
for i, (title_t, desc, col) in enumerate([
    ("Cuando usarla", "Despues de completar el trabajo, para confirmar con el cliente.", ACC),
    ("Quien firma", "Puede firmar el tecnico o el cliente.", GRN),
    ("Como firmar", "Usa tu dedo para dibujar directamente sobre el area.", PUR),
    ("Guardar", "Toca 'Guardar firma' para registrarla. No se edita despues.", RED),
    ("Limpiar", "Si te equivocas, toca 'Limpiar' para borrar y volver a firmar.", AMB),
    ("Confirmacion", "La firma queda registrada con fecha y hora automatica.", GRN),
]):
    y = 1.4 + i*0.85
    _rect(sl, R1X+0.2, y+0.05, 0.06, 0.7, col)
    _txt(sl, R1X+0.4, y+0.02, R1W-0.7, 0.2, title_t, 12, col, True)
    _txt(sl, R1X+0.4, y+0.28, R1W-0.7, 0.5, desc, 10, MUT)
print("Slide 9 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 10 - NOTAS
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Notas Internas", "Registra observaciones y comentarios sobre el reporte")
sx, sy, sw, sh = _phone(sl)
_topbar(sl, sx, sy, sw)
_txt(sl, sx+0.15, sy+0.52, 1.2, 0.22, "< Volver", 9, ACC, True)
_txt(sl, sx+0.15, sy+0.82, sw-0.3, 0.22, "NOTAS", 9, MUT, True)
n1 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(sy+1.12), Inches(sw-0.24), Inches(0.65))
n1.fill.solid(); n1.fill.fore_color.rgb = SURF2; n1.line.fill.background()
_txt(sl, sx+0.22, sy+1.17, sw-0.44, 0.3, "Sello mecanico danado. Necesito repuesto de la marca original.", 8, TXT2)
_txt(sl, sx+0.22, sy+1.5, sw-0.44, 0.16, "Carlos - 15/08 10:45", 7, MUT)
n2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(sy+1.87), Inches(sw-0.24), Inches(0.55))
n2.fill.solid(); n2.fill.fore_color.rgb = SURF2; n2.line.fill.background()
_txt(sl, sx+0.22, sy+1.92, sw-0.44, 0.25, "Cliente confirmo que el repuesto llega manana.", 8, TXT2)
_txt(sl, sx+0.22, sy+2.17, sw-0.44, 0.16, "Carlos - 15/08 14:20", 7, MUT)
_txt(sl, sx+0.15, sy+2.55, sw-0.3, 0.14, "NUEVA NOTA", 7, MUT, True)
_card(sl, sx+0.12, sy+2.72, sw-0.24, 0.65, BG2)
_txt(sl, sx+0.22, sy+2.78, sw-0.44, 0.5, "Escribe tu nota...", 9, MUT)
_nb = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(sy+3.48), Inches(sw-0.24), Inches(0.32))
_nb.fill.solid(); _nb.fill.fore_color.rgb = ACC; _nb.line.fill.background()
_txt(sl, sx+0.12, sy+3.5, sw-0.24, 0.28, "Enviar nota", 10, WHT, True, PP_ALIGN.CENTER)

_card(sl, R1X, 0.85, R1W, 5.9)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Uso de notas", 16, ACC, True)
for i, (title_t, desc, col) in enumerate([
    ("Que son", "Comentarios internos sobre el reporte. Solo el equipo los ve.", ACC),
    ("Para que sirven", "Registrar avances, hallazgos y decisiones durante la reparacion.", GRN),
    ("Agregar nota", "Escribe en el campo de texto y toca 'Enviar nota'.", PUR),
    ("Historial", "Las notas se ordenan cronologicamente y muestran autor y fecha.", AMB),
    ("Ejemplos", "Repuesto en camino, foto del dano, cambio de prioridad.", MUT),
]):
    y = 1.4 + i*0.95
    _rect(sl, R1X+0.2, y+0.05, 0.06, 0.8, col)
    _txt(sl, R1X+0.4, y+0.02, R1W-0.7, 0.2, title_t, 12, col, True)
    _txt(sl, R1X+0.4, y+0.28, R1W-0.7, 0.55, desc, 10, MUT)
print("Slide 10 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 11 - OFFLINE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Modo Offline", "Sigue trabajando aunque no tengas conexion a internet")
sx, sy, sw, sh = _phone(sl)
_topbar(sl, sx, sy, sw)
_ob = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(sx), Inches(sy+0.42), Inches(sw), Inches(0.24))
_ob.fill.solid(); _ob.fill.fore_color.rgb = AMB; _ob.line.fill.background()
_txt(sl, sx+0.1, sy+0.43, sw-0.2, 0.22, "Sin conexion - modo offline", 8, RGBColor(0x1a, 0x12, 0x00), True, PP_ALIGN.CENTER)
_txt(sl, sx+0.15, sy+0.78, sw-0.3, 0.24, "Mis reportes", 13, TXT, True)
_txt(sl, sx+0.15, sy+1.02, sw-0.3, 0.16, "Ultima sync: hace 15 min", 8, AMB)
_card(sl, sx+0.12, sy+1.28, sw-0.24, 0.4, AMB_BG)
_txt(sl, sx+0.22, sy+1.3, sw-0.44, 0.16, "2 reportes pendientes de sync", 8, AMB, True, PP_ALIGN.CENTER)
_txt(sl, sx+0.22, sy+1.48, sw-0.44, 0.16, "Se enviaran cuando haya internet", 7, MUT, False, PP_ALIGN.CENTER)
y5 = sy + 1.85
c5 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(y5), Inches(sw-0.24), Inches(0.85))
c5.fill.solid(); c5.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x30); c5.line.color.rgb = AMB; c5.line.width = Pt(0.5)
_rect(sl, sx+0.12, y5, 0.03, 0.85, AMB)
_txt(sl, sx+0.24, y5+0.06, 1.8, 0.2, "Local - Pendiente", 9, AMB, True)
_txt(sl, sx+0.24, y5+0.28, sw-0.44, 0.16, "Nueva empresa S.A.", 8, TXT, True)
_txt(sl, sx+0.24, y5+0.48, sw-0.44, 0.3, "Bomba de agua con ruido anormal en el rodamiento", 8, TXT2)

_card(sl, R1X, 0.85, R1W, 3.1)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Flujo offline", 16, AMB, True)
for i, (num, title_t, desc, col) in enumerate([
    ("1", "Sin internet", "La app detecta que no hay conexion", RED),
    ("2", "Guarda local", "El reporte se guarda en el dispositivo", AMB),
    ("3", "Cola pendiente", "Aparece indicador de pendientes", PUR),
    ("4", "Internet vuelve", "La app detecta la conexion", GRN),
    ("5", "Sync automatica", "Todos los pendientes se envian", ACC),
]):
    y = 1.4 + i*0.48
    _circ(sl, R1X+0.2, y+0.02, 0.3, col)
    _txt(sl, R1X+0.2, y+0.04, 0.3, 0.22, num, 9, WHT, True, PP_ALIGN.CENTER)
    _txt(sl, R1X+0.6, y, 1.4, 0.16, title_t, 10, col, True)
    _txt(sl, R1X+2.05, y, 1.75, 0.16, desc, 9, MUT)
    if i < 4:
        _rect(sl, R1X+0.33, y+0.34, 0.04, 0.12, MUT)

_card(sl, R1X, 4.1, R1W, 2.6)
_txt(sl, R1X+0.2, 4.2, R1W-0.4, 0.3, "Que funciona offline", 14, GRN, True)
for i, (item, col) in enumerate([
    ("Crear reportes", GRN), ("Ver reportes locales", GRN), ("Tomar fotos", GRN),
    ("Firma digital", GRN), ("Filtros de busqueda", GRN), ("Notas internas", GRN),
]):
    y = 4.6 + i*0.32
    _circ(sl, R1X+0.2, y+0.03, 0.14, col)
    _txt(sl, R1X+0.45, y, 3.0, 0.2, item, 11, TXT, True)

_card(sl, R2X, 0.85, R2W, 5.9)
_txt(sl, R2X+0.2, 0.95, R2W-0.4, 0.3, "Detalles tecnicos", 16, PUR, True)
for i, (item, desc, col) in enumerate([
    ("IndexedDB", "Los datos se guardan en la base de datos del navegador.", ACC),
    ("Sync cada 30s", "La app revisa la conexion cada 30 segundos.", GRN),
    ("Fotos comprimidas", "Las imagenes se reducen de tamano antes de enviarse.", AMB),
    ("Cola persistente", "Los datos sobreviven cierre de app o reinicio.", RED),
    ("Conflictos", "Si hay duplicados, el servidor resuelve por timestamp.", PUR),
    ("Sin limite", "No hay limite de reportes offline. Suben cuando haya red.", MUT),
]):
    y = 1.4 + i*0.82
    _rect(sl, R2X+0.2, y+0.05, 0.06, 0.65, col)
    _txt(sl, R2X+0.4, y+0.02, R2W-0.7, 0.2, item, 11, col, True)
    _txt(sl, R2X+0.4, y+0.28, R2W-0.7, 0.45, desc, 10, MUT)
print("Slide 11 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 12 - DASHBOARD
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Estadisticas (Dashboard)", "Resumen de tu actividad y metricas de rendimiento")
sx, sy, sw, sh = _phone(sl)
_topbar(sl, sx, sy, sw)
_txt(sl, sx+0.15, sy+0.52, 1.2, 0.22, "< Volver", 9, ACC, True)
_txt(sl, sx+0.15, sy+0.82, sw-0.3, 0.24, "Estadisticas", 12, TXT, True)
for i, (num, lbl, col) in enumerate([("24", "Total", ACC), ("8", "Abiertos", RED), ("6", "En pro", ACC2), ("7", "Resueltos", GRN), ("3", "Espera", AMB), ("3d", "Promedio", ACC)]):
    r = i // 3; c = i % 3
    cx = sx + 0.12 + c*(sw-0.24)/3
    cy = sy + 1.15 + r*0.95
    cw = (sw-0.24)/3 - 0.04
    sc = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cx), Inches(cy), Inches(cw), Inches(0.85))
    sc.fill.solid(); sc.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x30); sc.line.color.rgb = col; sc.line.width = Pt(0.5)
    _txt(sl, cx, cy+0.1, cw, 0.35, num, 18, col, True, PP_ALIGN.CENTER)
    _txt(sl, cx, cy+0.52, cw, 0.18, lbl, 7, MUT, False, PP_ALIGN.CENTER)
_txt(sl, sx+0.15, sy+3.15, sw-0.3, 0.14, "POR PRIORIDAD", 7, MUT, True)
for i, (lbl, pct, col) in enumerate([("Urg", 30, RED), ("Alt", 45, ORG), ("Nor", 70, ACC), ("Baj", 20, GGRAY)]):
    by = sy + 3.4 + i*0.28
    _txt(sl, sx+0.15, by, 0.35, 0.18, lbl, 7, MUT)
    _rect(sl, sx+0.52, by+0.03, (sw-0.76)*pct/100, 0.12, col)
    _rect(sl, sx+0.52+(sw-0.76)*pct/100, by+0.03, (sw-0.76)*(100-pct)/100, 0.12, SURF2)

_card(sl, R1X, 0.85, R1W, 5.9)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Acceso y uso", 16, ACC, True)
for i, (title_t, desc, col) in enumerate([
    ("Como acceder", "Toca el menu (...) en la parte superior derecha, luego 'Estadisticas'.", ACC),
    ("Stat cards", "Muestra 6 metricas clave: Total, Abiertos, En proceso, Resueltos, En espera, Promedio.", GRN),
    ("Barras por prioridad", "Grafico de barras con reportes por nivel de prioridad.", RED),
    ("Ultimos 7 dias", "Cuantos reportes se crearon en la ultima semana.", AMB),
    ("Datos en tiempo real", "Las estadisticas se actualizan cada vez que accedes.", PUR),
]):
    y = 1.4 + i*0.95
    _rect(sl, R1X+0.2, y+0.05, 0.06, 0.8, col)
    _txt(sl, R1X+0.4, y+0.02, R1W-0.7, 0.2, title_t, 12, col, True)
    _txt(sl, R1X+0.4, y+0.28, R1W-0.7, 0.55, desc, 10, MUT)
print("Slide 12 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 13 - ARCHIVADOS
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Reportes Archivados", "Reportes resueltos automaticamente despues de 24 horas")
sx, sy, sw, sh = _phone(sl)
_topbar(sl, sx, sy, sw)
_txt(sl, sx+0.15, sy+0.52, 1.2, 0.22, "< Volver", 9, ACC, True)
_txt(sl, sx+0.15, sy+0.82, sw-0.3, 0.24, "Archivados (5)", 12, TXT, True)
_txt(sl, sx+0.15, sy+1.06, sw-0.3, 0.16, "Resueltos con mas de 24h", 8, MUT)
for i, (cli, eq, tech) in enumerate([
    ("Ecopetrol SA", "Turbina TG-200", "Carlos"),
    ("Cementos Argos", "Horno rotary HR-3", "Ana"),
    ("Minera del Sur", "Trituradora TR-50", "Luis"),
]):
    ay = sy + 1.3 + i*0.9
    ac = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(sx+0.12), Inches(ay), Inches(sw-0.24), Inches(0.78))
    ac.fill.solid(); ac.fill.fore_color.rgb = RGBColor(0x11, 0x18, 0x30); ac.line.color.rgb = RGBColor(0x2a, 0x3a, 0x5a); ac.line.width = Pt(0.3)
    _rect(sl, sx+0.12, ay, 0.03, 0.78, GRN)
    _txt(sl, sx+0.24, ay+0.06, 1.8, 0.18, cli, 9, TXT, True)
    _mbadge(sl, sx+sw-0.75, ay+0.06, "Resuelto", WHT, GRN_BG)
    _txt(sl, sx+0.24, ay+0.28, sw-0.44, 0.14, eq, 8, MUT)
    _txt(sl, sx+0.24, ay+0.48, 1.5, 0.14, tech, 7, MUT)

_card(sl, R1X, 0.85, R1W, 3.2)
_txt(sl, R1X+0.2, 0.95, R1W-0.4, 0.3, "Sobre los archivados", 16, GRN, True)
for i, (title_t, desc, col) in enumerate([
    ("Auto-archivado", "Los reportes en estado 'resuelto' se archivan despues de 24 horas.", GRN),
    ("Acceder", "Menu (...) > Archivados. Se muestra la lista completa.", ACC),
    ("Solo lectura", "Los archivados son de consulta. No se pueden editar.", MUT),
    ("Detalle completo", "Toca un archivado para ver toda su informacion original.", PUR),
]):
    y = 1.4 + i*0.65
    _rect(sl, R1X+0.2, y+0.05, 0.06, 0.5, col)
    _txt(sl, R1X+0.4, y+0.02, R1W-0.7, 0.2, title_t, 11, col, True)
    _txt(sl, R1X+0.4, y+0.22, R1W-0.7, 0.4, desc, 10, MUT)

_card(sl, R1X, 4.2, R1W, 2.5)
_txt(sl, R1X+0.2, 4.3, R1W-0.4, 0.3, "Menu desplegable", 14, ACC, True)
for i, (item, desc, col) in enumerate([
    ("Monitoreo", "Ver todos los reportes (solo gerente)", ACC),
    ("Archivados", "Reportes resueltos +24h", GRN),
    ("Estadisticas", "Dashboard con metricas", PUR),
    ("Actualizar", "Sync manual de datos", AMB),
    ("Tema", "Cambiar claro/oscuro", MUT),
    ("Cerrar sesion", "Salir de la app", RED),
]):
    y = 4.7 + i*0.32
    _txt(sl, R1X+0.3, y, 1.5, 0.22, item, 10, col, True)
    _txt(sl, R1X+1.9, y, 1.9, 0.22, desc, 9, MUT)
print("Slide 13 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 14 - PUSH
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
_title(sl, "Notificaciones Push", "Recibe alertas importantes en tu dispositivo")
_card(sl, 0.5, 0.9, 5.8, 2.1)
_txt(sl, 0.7, 1.0, 5.4, 0.3, "Ejemplo de notificaciones", 15, ACC, True)
for i, (title_t, body, col) in enumerate([
    ("NexAlert", "Se te asigno el reporte #45 de Compresores XYZ por falla en motor", ACC),
    ("NexAlert - URGENTE", "Nuevo reporte urgente: Parada total de linea de produccion en Minera del Sur", RED),
    ("NexAlert", "El reporte #23 lleva 8 dias abierto. Por favor revisar y actualizar.", AMB),
]):
    y = 1.45 + i*0.5
    nc = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(y), Inches(5.4), Inches(0.45))
    nc.fill.solid(); nc.fill.fore_color.rgb = SURF2; nc.line.fill.background()
    _rect(sl, 0.7, y, 0.04, 0.45, col)
    _txt(sl, 0.88, y+0.02, 5.0, 0.18, title_t, 9, col, True)
    _txt(sl, 0.88, y+0.2, 5.0, 0.22, body, 8, MUT)

_card(sl, 0.5, 3.2, 5.8, 3.5)
_txt(sl, 0.7, 3.3, 5.4, 0.3, "Eventos que generan push", 15, ACC, True)
for i, (ev, desc, col) in enumerate([
    ("Asignacion", "Te asignan un reporte nuevo", ACC),
    ("Cambio de estado", "El estado de tu reporte cambia", GRN),
    ("Nuevo mensaje", "El cliente responde por WhatsApp", AMB),
    ("Urgente", "Se crea un reporte urgente", RED),
    ("Recordatorio", "Reporte sin resolver por +7 dias", ORG),
]):
    y = 3.7 + i*0.58
    _rect(sl, 0.7, y+0.05, 0.06, 0.42, col)
    _txt(sl, 0.9, y+0.02, 2.2, 0.2, ev, 11, col, True)
    _txt(sl, 3.2, y+0.02, 2.8, 0.2, desc, 10, MUT)

_card(sl, 6.6, 0.9, 6.2, 5.8)
_txt(sl, 6.8, 1.0, 5.8, 0.3, "Configuracion de notificaciones", 16, ACC, True)
for i, (title_t, desc, col) in enumerate([
    ("Acepta los permisos", "Al instalar la app, acepta recibir notificaciones push.", ACC),
    ("Funciona con app cerrada", "Las notificaciones llegan incluso si la app no esta abierta.", GRN),
    ("Toca para abrir", "Al tocar la notificacion, se abre la app en el reporte relevante.", PUR),
    ("No desactives las push", "Son esenciales para no perderte asignaciones y cambios importantes.", RED),
    ("Firebase", "Usamos Firebase Cloud Messaging para entregar notificaciones confiables.", MUT),
    ("Token unico", "Cada dispositivo tiene un token unico para recibir sus notificaciones.", AMB),
    ("Actualizaciones", "Tambien recibes aviso cuando hay una nueva version de la app.", GRN),
]):
    y = 1.4 + i*0.7
    _rect(sl, 6.8, y+0.05, 0.06, 0.55, col)
    _txt(sl, 7.0, y+0.02, 5.4, 0.2, title_t, 12, col, True)
    _txt(sl, 7.0, y+0.25, 5.4, 0.4, desc, 10, MUT)
print("Slide 14 OK")

# ═══════════════════════════════════════════════════════════
# SLIDE 15 - CIERRE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); _slide(sl)
ov = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
ov.fill.solid(); ov.fill.fore_color.rgb = RGBColor(0x06, 0x0a, 0x14); ov.line.fill.background()
_circ(sl, 6.15, 0.4, 1.1, ACC)
_txt(sl, 6.15, 0.52, 1.1, 0.7, "N", 34, WHT, True, PP_ALIGN.CENTER)
_txt(sl, 1, 1.8, 11.3, 0.7, "GRACIAS POR SU ATENCION", 40, WHT, True, PP_ALIGN.CENTER)
_txt(sl, 1, 2.7, 11.3, 0.4, "Estas listo para usar NexAlert en campo", 18, MUT, False, PP_ALIGN.CENTER)
d = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(3.3), Inches(2.7), Inches(0.03))
d.fill.solid(); d.fill.fore_color.rgb = ACC; d.line.fill.background()
_card(sl, 2.0, 3.7, 9.3, 3.1)
_txt(sl, 2.2, 3.8, 8.9, 0.3, "Lo que aprendiste hoy", 18, ACC, True)
for i, (item, desc, col) in enumerate([
    ("Login y biometria", "Acceder a la app de forma segura", ACC),
    ("Crear reportes", "Documentar fallas con descripcion clara", GRN),
    ("Fotografias", "Capturar evidencia visual del problema", PUR),
    ("Cambiar estados", "Actualizar el progreso de la reparacion", AMB),
    ("Firma digital", "Confirmar entrega con firma del cliente", GRN),
    ("Modo offline", "Trabajar sin conexion y sincronizar despues", RED),
    ("Dashboard", "Ver estadisticas y metricas de rendimiento", ACC),
    ("Archivados", "Consultar reportes resueltos automaticamente", MUT),
]):
    x = 2.3 + (i%2)*4.6
    y = 4.25 + (i//2)*0.55
    _circ(sl, x, y+0.03, 0.2, col)
    _txt(sl, x, y+0.02, 0.2, 0.16, "V", 8, WHT, True, PP_ALIGN.CENTER)
    _txt(sl, x+0.3, y, 2.0, 0.18, item, 12, TXT, True)
    _txt(sl, x+0.3, y+0.2, 4.0, 0.18, desc, 10, MUT)
_txt(sl, 1, 7.0, 11.3, 0.3, "NEXUS  |  v1.5.5  |  2026", 12, MUT, False, PP_ALIGN.CENTER)
print("Slide 15 OK")

prs.save(r"C:\Users\STIVEN\Desktop\NexAlert-Capacitacion.pptx")
print(f"\nGuardado: {len(prs.slides)} slides")
