from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors matching actual CSS
BG = RGBColor(0x0a, 0x0e, 0x1a)
SURFACE = RGBColor(0x14, 0x1b, 0x30)
SURFACE2 = RGBColor(0x1b, 0x24, 0x40)
ACCENT = RGBColor(0x4f, 0x7d, 0xff)
ACCENT2 = RGBColor(0x7a, 0x9d, 0xff)
GREEN = RGBColor(0x2f, 0xd0, 0x7f)
RED = RGBColor(0xff, 0x5c, 0x6c)
AMBER = RGBColor(0xff, 0xb4, 0x54)
ORANGE = RGBColor(0xff, 0x8a, 0x4c)
TEXT = RGBColor(0xee, 0xf2, 0xff)
TEXT2 = RGBColor(0xd4, 0xdb, 0xf2)
MUTED = RGBColor(0x8b, 0x95, 0xb8)
BORDER = RGBColor(0x2a, 0x3a, 0x5a)
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK = RGBColor(0x1a, 0x1c, 0x23)
BLUE_BG = RGBColor(0x1a, 0x20, 0x3a)
RED_BG = RGBColor(0x3a, 0x15, 0x1c)
GREEN_BG = RGBColor(0x0f, 0x2a, 0x1c)
AMBER_BG = RGBColor(0x3a, 0x2e, 0x15)

SCREENSHOTS = r"C:\Users\STIVEN\Documents\Default Project\nexalert-screenshots"
LOGO = r"C:\Users\STIVEN\Documents\Default Project\reportes-equipos\build\icon-512.png"

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def topbar(slide, color=BG):
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    bbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.42), Inches(13.333), Inches(0.08))
    bbar.fill.solid(); bbar.fill.fore_color.rgb = ACCENT; bbar.line.fill.background()

def title(slide, text, sub=None):
    box = s.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(12), Inches(0.5))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(28); p.font.color.rgb = TEXT; p.font.bold = True
    if sub:
        p2 = tf.add_paragraph(); p2.text = sub; p2.font.size = Pt(14); p2.font.color.rgb = MUTED

def txt(slide, l, t, w, h, text, size=14, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size); p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tf

def bullets(slide, l, t, w, h, items, size=13, color=TEXT, spacing=6):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(size); p.font.color.rgb = color; p.space_after = Pt(spacing)
    return tf

def phone_mockup(slide, left, top, width, height):
    body = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    body.fill.solid(); body.fill.fore_color.rgb = RGBColor(0x0f, 0x15, 0x26); body.line.fill.background()
    notch = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + width*0.35), Inches(top), Inches(width*0.3), Inches(0.12))
    notch.fill.solid(); notch.fill.fore_color.rgb = RGBColor(0x0a, 0x0e, 0x1a); notch.line.fill.background()
    inner = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left+0.06), Inches(top+0.1), Inches(width-0.12), Inches(height-0.16))
    inner.fill.solid(); inner.fill.fore_color.rgb = BG; inner.line.fill.background()
    return inner

def monitor_mockup(slide, left, top, width, height):
    stand = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + width/2 - 0.15), Inches(top + height), Inches(0.3), Inches(0.15))
    stand.fill.solid(); stand.fill.fore_color.rgb = RGBColor(0x3a, 0x3a, 0x3a); stand.line.fill.background()
    base = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + width/2 - 0.35), Inches(top + height + 0.15), Inches(0.7), Inches(0.05))
    base.fill.solid(); base.fill.fore_color.rgb = RGBColor(0x3a, 0x3a, 0x3a); base.line.fill.background()
    frame = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    frame.fill.solid(); frame.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x2d); frame.line.fill.background()
    inner = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left+0.05), Inches(top+0.05), Inches(width-0.1), Inches(height-0.1))
    inner.fill.solid(); inner.fill.fore_color.rgb = BG; inner.line.fill.background()
    return inner

def pill(slide, l, t, text, color, bg_color, w=None):
    w = w or max(0.8, len(text)*0.12)
    shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.22))
    shape.fill.solid(); shape.fill.fore_color.rgb = bg_color; shape.line.fill.background()
    txt(slide, l, t+0.01, w, 0.2, text, 9, color, True, PP_ALIGN.CENTER)

def stat_card(slide, l, t, num, label, num_color):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(1.1), Inches(0.85))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.color.rgb = BORDER; card.line.width = Pt(0.5)
    txt(slide, l, t+0.08, 1.1, 0.4, str(num), 22, num_color, True, PP_ALIGN.CENTER)
    txt(slide, l, t+0.5, 1.1, 0.3, label, 9, MUTED, False, PP_ALIGN.CENTER)

def mockup_list_item(slide, l, t, w, name, detail, badge_text, badge_color, badge_bg):
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.55))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.color.rgb = BORDER; card.line.width = Pt(0.3)
    txt(slide, l+0.1, t+0.05, w*0.55, 0.25, name, 11, TEXT, True)
    txt(slide, l+0.1, t+0.28, w*0.55, 0.22, detail, 9, MUTED)
    pill(slide, l+w-0.9, t+0.16, badge_text, badge_color, badge_bg, 0.85)

# ===== SLIDE 1: PORTADA =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK)
topbar(s)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.9), Inches(1.0), Inches(1.5), Inches(1.5))
txt(s, 1, 2.8, 11.3, 0.9, "CAPACITACION NEXALERT", 44, WHITE, True, PP_ALIGN.CENTER)
txt(s, 1, 3.8, 11.3, 0.6, "Sistema integral de gestion de reportes de fallas", 18, MUTED, False, PP_ALIGN.CENTER)
txt(s, 1, 4.6, 11.3, 0.6, "Desktop  |  App Movil  |  Servidor  |  WhatsApp", 16, ACCENT, False, PP_ALIGN.CENTER)
txt(s, 1, 5.8, 11.3, 0.5, "NEXUS  |  2026  |  v1.5.5", 13, MUTED, False, PP_ALIGN.CENTER)

# ===== SLIDE 2: QUE ES NEXALERT =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Que es NexAlert?", "Plataforma completa para gestionar reportes de fallas de equipos")

# Desktop mockup
inner = monitor_mockup(s, 0.5, 1.5, 5.5, 3.5)
# Topbar in monitor
tb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.55), Inches(5.4), Inches(0.35))
tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0x1b, 0x1f, 0x3b); tb.line.fill.background()
# Tab bar
tabbar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.9), Inches(5.4), Inches(0.28))
tabbar.fill.solid(); tabbar.fill.fore_color.rgb = SURFACE; tabbar.line.fill.background()
tabs = ["Panel", "Clientes", "Equipos", "Reportes", "Mensajes", "Tecnicos"]
for i, tab in enumerate(tabs):
    c = ACCENT if tab == "Reportes" else MUTED
    txt(s, 0.6 + i*0.9, 1.92, 0.85, 0.24, tab, 7, c, tab == "Reportes")
# Stats
for i, (num, label, col) in enumerate([(47,"Reportes",ACCENT),(7,"Abiertos",RED),(6,"Proceso",AMBER),(24,"Resueltos",GREEN)]):
    stat_card(s, 0.65 + i*1.3, 2.3, num, label, col)

txt(s, 6.5, 1.5, 6.5, 0.4, "Componentes del sistema:", 16, TEXT, True)
items = [
    "Desktop (Electron) - Panel de gerente con 6 modulos",
    "App movil (Android/Capacitor) - Para tecnicos en campo",
    "Servidor (Node.js + SQLite) - GCP us-central1",
    "WhatsApp Web - Integracion directa con clientes",
    "Firebase Push - Notificaciones en tiempo real",
    "Sincronizacion bidireccional - Online/Offline",
    "Auto-update - Actualizaciones automaticas",
    "Biometria - Login rapido con huella/rostro"
]
bullets(s, 6.5, 2.1, 6.3, 5, items, 12, TEXT2)

# Mobile mockup
inner2 = phone_mockup(s, 2, 5.3, 2.2, 2)
txt(s, 2.15, 5.45, 1.9, 0.2, "NexAlert", 9, ACCENT, True, PP_ALIGN.CENTER)
txt(s, 2.15, 5.7, 1.9, 0.15, "Mis reportes", 7, TEXT, True, PP_ALIGN.CENTER)
for i, (name, eq, estado, ec) in enumerate([
    ("tacos El Güero","Freidora #3","Abierto",RED),
    ("Cafe Azul","Computo #1","En proceso",AMBER),
    ("Farmacia Lopez","Refrigerador","Resuelto",GREEN)]):
    y = 5.95 + i*0.4
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.15), Inches(y), Inches(1.9), Inches(0.35))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.2), Inches(y+0.08), Inches(0.1), Inches(0.1))
    dot.fill.solid(); dot.fill.fore_color.rgb = ec; dot.line.fill.background()
    txt(s, 2.35, y+0.02, 1.1, 0.15, name, 6, TEXT, True)
    txt(s, 2.35, y+0.17, 1.1, 0.12, eq, 5, MUTED)

# ===== SLIDE 3: ARQUITECTURA =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Arquitectura del sistema", "Como se conectan los componentes")

# Desktop
inner = monitor_mockup(s, 0.3, 1.8, 4.2, 2.8)
txt(s, 0.3, 4.7, 4.2, 0.35, "Desktop (Electron)", 13, ACCENT, True, PP_ALIGN.CENTER)
txt(s, 0.3, 5.05, 4.2, 0.3, "Gerente / Oficina", 11, MUTED, False, PP_ALIGN.CENTER)

# Arrow right
arr1 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.7), Inches(3), Inches(0.8), Inches(0.3))
arr1.fill.solid(); arr1.fill.fore_color.rgb = ACCENT; arr1.line.fill.background()
# Arrow left (sync)
arr2 = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(4.7), Inches(3.5), Inches(0.8), Inches(0.3))
arr2.fill.solid(); arr2.fill.fore_color.rgb = GREEN; arr2.line.fill.background()
txt(s, 4.5, 3.9, 1.2, 0.2, "Sync", 8, MUTED, False, PP_ALIGN.CENTER)

# Server box
srv = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.8), Inches(1.8), Inches(2.3), Inches(3))
srv.fill.solid(); srv.fill.fore_color.rgb = SURFACE; srv.line.color.rgb = ACCENT; srv.line.width = Pt(1.5)
txt(s, 5.8, 1.9, 2.3, 0.3, "Servidor", 13, ACCENT, True, PP_ALIGN.CENTER)
server_items = [
    "Node.js + Express",
    "SQLite (WAL mode)",
    "Firebase Admin",
    "JWT Auth",
    "Seq Log (sync)",
    "Auto-archive +24h",
    "GCP us-central1",
    "HTTPS + DuckDNS"
]
bullets(s, 5.9, 2.3, 2.1, 2.3, server_items, 9, TEXT2, 3)

# Arrow right
arr3 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.3), Inches(3), Inches(0.8), Inches(0.3))
arr3.fill.solid(); arr3.fill.fore_color.rgb = ACCENT; arr3.line.fill.background()
arr4 = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(8.3), Inches(3.5), Inches(0.8), Inches(0.3))
arr4.fill.solid(); arr4.fill.fore_color.rgb = GREEN; arr4.line.fill.background()

# Mobile
inner2 = phone_mockup(s, 9.5, 1.8, 2.5, 2.8)
txt(s, 9.5, 4.7, 2.5, 0.35, "App Movil (Android)", 13, ACCENT, True, PP_ALIGN.CENTER)
txt(s, 9.5, 5.05, 2.5, 0.3, "Tecnico en campo", 11, MUTED, False, PP_ALIGN.CENTER)

# Features below
feat_cards = [
    ("Offline", "IndexedDB local\nCola de cambios", AMBER),
    ("WhatsApp", "QR Connect\nChat + Grupos", GREEN),
    ("Push", "Firebase FCM\nNotificaciones", RED),
    ("Biometria", "NativeBiometric\nLogin rapido", ACCENT),
    ("GPS", "Geolocation\nUbicacion exacta", GREEN),
]
for i, (name, desc, col) in enumerate(feat_cards):
    left = 0.3 + i * 2.6
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(5.6), Inches(2.3), Inches(1.3))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.9), Inches(5.7), Inches(0.5), Inches(0.5))
    dot.fill.solid(); dot.fill.fore_color.rgb = col; dot.line.fill.background()
    txt(s, left + 0.1, 6.25, 2.1, 0.25, name, 12, col, True, PP_ALIGN.CENTER)
    txt(s, left + 0.1, 6.5, 2.1, 0.4, desc, 9, MUTED, False, PP_ALIGN.CENTER)

# ===== SLIDE 4: ROLES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Quien usa que?", "Cada rol tiene acceso a diferentes funciones")

# Gerente card
g_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.3), Inches(6.1), Inches(5.5))
g_card.fill.solid(); g_card.fill.fore_color.rgb = SURFACE; g_card.line.color.rgb = ACCENT; g_card.line.width = Pt(1)
txt(s, 0.6, 1.4, 5.7, 0.4, "Gerente (Desktop + Movil)", 18, ACCENT, True)
txt(s, 0.6, 1.9, 5.7, 0.3, "Desktop - Panel de control completo:", 13, TEXT, True)
bullets(s, 0.8, 2.2, 5.3, 2, [
    "Panel: 7 tarjetas estadisticas + 7 paneles de analisis",
    "Clientes: CRUD + grupos de WhatsApp",
    "Equipos: CRUD por cliente + historial de fallas",
    "Reportes: 6 filtros + busqueda + exportar PDF/Excel",
    "Mensajes: WhatsApp Web completo (chat, emoji, stickers, menciones)",
    "Tecnicos: CRUD + importar desde archivo + activar/desactivar"
], 11, TEXT2, 4)
txt(s, 0.6, 4.4, 5.7, 0.3, "Movil - Monitoreo remoto:", 13, TEXT, True)
bullets(s, 0.8, 4.7, 5.3, 1.5, [
    "Monitorear tecnicos en tiempo real",
    "Ver todos los reportes + estadisticas globales",
    "Asignar/reasignar tecnicos desde el celular",
    "Cambiar estados + desactivar tecnicos"
], 11, TEXT2, 4)

# Tecnico card
t_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.3), Inches(6.1), Inches(5.5))
t_card.fill.solid(); t_card.fill.fore_color.rgb = SURFACE; t_card.line.color.rgb = GREEN; t_card.line.width = Pt(1)
txt(s, 7, 1.4, 5.7, 0.4, "Tecnico (Solo Movil)", 18, GREEN, True)
bullets(s, 7.2, 1.9, 5.3, 4.5, [
    "Ver reportes asignados (solo los suyos)",
    "Crear reportes con: cliente, equipo, descripcion, prioridad, foto, GPS",
    "Cambiar estado: Abierto / En proceso / Espera repuesto / Espera cliente / Resuelto",
    "Firma del cliente al resolver (canvas touch)",
    "Adjuntar fotos: camara o galeria (compress automatico 1280px)",
    "Compartir reporte por Capacitor Share",
    "Trabajar offline: cola de cambios en IndexedDB",
    "Biometria para login rapido",
    "Recibir notificaciones push",
    "Dashboard de estadisticas personales",
    "Ver archivados (reportes resueltos +24h)"
], 11, TEXT2, 4)

# ===== SLIDE 5: DESKTOP - PANEL =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Desktop: Panel Principal", "Centro de control del gerente")

# Stats row
stats = [
    (7, "Abiertos", RED), (6, "En proceso", AMBER), (4, "En espera", MUTED),
    (3, "Urgentes", ORANGE), (2, "Sin enviar", ACCENT), (5, "Sin tecnico", MUTED), (24, "Resueltos", GREEN)
]
for i, (n, l, c) in enumerate(stats):
    stat_card(s, 0.4 + i*1.8, 1.3, n, l, c)

# Panel cards
panel_items = [
    ("Carga por tecnico", ["Juan: 3 en curso, 5 resueltos", "Pedro: 2 en curso, 8 resueltos", "Maria: 1 en curso, 4 resueltos", "Sin asignar: 5 reportes"]),
    ("Reportes activos por cliente", ["tacos El Güero: 4", "Cafe Azul: 3", "Farmacia Lopez: 2", "AutoZone: 1"]),
    ("Reportes urgentes", ["tacos El Güero - Freidora #3", "Cafe Azul - Computo #1", "Farmacia Lopez - Refrigerador"]),
    ("Equipos que mas fallan", ["Freidora #3: 8 fallas", "Computo #1: 5 fallas", "Refrigerador: 4 fallas"]),
    ("Promedio dias resolver", ["Global: 2.3 dias", "Juan: 1.8 dias", "Pedro: 2.9 dias"]),
    ("Antiguos +7 dias", ["tacos El Güero - Freidora (12d)", "Cafe Azul - Computo (9d)"]),
]
for i, (title_txt, items_list) in enumerate(panel_items):
    col = i % 3
    row = i // 3
    left = 0.4 + col * 4.2
    top = 2.5 + row * 2.4
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.9), Inches(2.15))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    txt(s, left+0.15, top+0.08, 3.6, 0.25, title_txt, 11, TEXT, True)
    for j, item in enumerate(items_list):
        txt(s, left+0.15, top+0.35 + j*0.35, 3.6, 0.3, item, 9, TEXT2)

# ===== SLIDE 6: DESKTOP - REPORTES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Desktop: Gestion de Reportes", "Filtros, busqueda y acciones completas")

# Filter bar
fbar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.2), Inches(12.7), Inches(0.65))
fbar.fill.solid(); fbar.fill.fore_color.rgb = SURFACE; fbar.line.fill.background()

filters = ["Todos los reportes", "Urgentes", "Abiertos", "En proceso", "En espera", "Resueltos", "Pendientes", "Sin tecnico", "Archivados"]
for i, f in enumerate(filters):
    c = ACCENT if f == "Todos los reportes" else MUTED
    txt(s, 0.4 + i*1.38, 1.32, 1.35, 0.2, f, 8, c, f == "Todos los reportes")
txt(s, 0.4, 1.52, 2, 0.2, "Buscar por texto... | Este mes | Todos los clientes", 7, MUTED)
# Buttons
btn1 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10), Inches(1.28), Inches(1.3), Inches(0.3))
btn1.fill.solid(); btn1.fill.fore_color.rgb = ACCENT; btn1.line.fill.background()
txt(s, 10, 1.29, 1.3, 0.28, "+ Nuevo reporte", 9, WHITE, True, PP_ALIGN.CENTER)
btn2 = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.4), Inches(1.28), Inches(0.8), Inches(0.3))
btn2.fill.solid(); btn2.fill.fore_color.rgb = SURFACE2; btn2.line.fill.background()
txt(s, 11.4, 1.29, 0.8, 0.28, "Excel", 9, TEXT2, False, PP_ALIGN.CENTER)

# Report list items
reports = [
    ("tacos El Güero", "Freidora #3", "No enciende, se escucha ruido raro", "abierto", "urgente", "Juan Perez"),
    ("Cafe Azul", "Computo #1", "Pantalla azul despues de actualizacion Windows", "en_proceso", "alta", "Pedro Lopez"),
    ("Farmacia Lopez", "Refrigerador #2", "No enfría, compresor no enciende", "espera_repuesto", "normal", "Maria Garcia"),
    ("AutoZone", "Impresora Termica", "Imprime con lineas, cabezal sucio", "resuelto", "baja", "Juan Perez"),
]
for i, (client, equipo, desc, estado, prio, tech) in enumerate(reports):
    y = 2.1 + i*0.9
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(y), Inches(9.5), Inches(0.78))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    # Left accent
    left_accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(y), Inches(0.05), Inches(0.78))
    left_accent.fill.solid()
    if estado == "abierto": left_accent.fill.fore_color.rgb = RED
    elif estado == "en_proceso": left_accent.fill.fore_color.rgb = AMBER
    elif estado == "resuelto": left_accent.fill.fore_color.rgb = GREEN
    else: left_accent.fill.fore_color.rgb = MUTED
    left_accent.line.fill.background()
    txt(s, 0.5, y+0.05, 2, 0.2, client, 12, TEXT, True)
    txt(s, 0.5, y+0.25, 2, 0.18, equipo, 9, MUTED)
    txt(s, 2.7, y+0.08, 5, 0.3, desc, 10, TEXT2)
    pill(s, 0.5, y+0.5, estado.replace("_"," ").title(), TEXT, SURFACE2 if estado=="resuelto" else RED_BG if estado=="abierto" else AMBER_BG)
    pill(s, 2.2, y+0.5, prio.title(), WHITE, RED if prio=="urgente" else ORANGE if prio=="alta" else ACCENT if prio=="normal" else MUTED)
    txt(s, 5, y+0.55, 2, 0.18, tech, 8, MUTED)

# Action buttons
actions = ["Cambiar estado", "Asignar", "Compartir", "PDF", "WhatsApp", "Eliminar"]
for i, act in enumerate(actions):
    btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.1), Inches(2.1 + i*0.55), Inches(2.8), Inches(0.42))
    btn.fill.solid(); btn.fill.fore_color.rgb = SURFACE2; btn.line.fill.background()
    txt(s, 10.2, y+0.06 if i==0 else 2.1 + i*0.55 + 0.06, 2.6, 0.3, act, 10, TEXT2)

txt(s, 10.1, 1.8, 2.8, 0.25, "Acciones del reporte:", 11, TEXT, True)

# ===== SLIDE 7: DESKTOP - CREAR REPORTE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Desktop: Crear Reporte", "Desde el boton '+ Nuevo reporte' se abre un modal")

# Modal mockup
modal = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.3), Inches(7), Inches(5.8))
modal.fill.solid(); modal.fill.fore_color.rgb = WHITE; modal.line.fill.background()
# Header
modal_head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.3), Inches(7), Inches(0.5))
modal_head.fill.solid(); modal_head.fill.fore_color.rgb = RGBColor(0xf8, 0xfa, 0xfc); modal_head.line.fill.background()
txt(s, 0.8, 1.35, 5, 0.4, "Nuevo reporte de falla", 16, DARK, True)
txt(s, 7, 1.38, 0.4, 0.35, "X", 14, MUTED, True, PP_ALIGN.CENTER)

fields = [
    ("Cliente *", "Selecciona o escribe el nombre", True),
    ("Equipo", "Marca, modelo, serie...", False),
    ("Descripcion del problema *", "Describe el problema detalladamente...", False),
    ("Prioridad", "Normal  (baja | normal | alta | urgente)", False),
    ("Grupo de WhatsApp", "Selecciona el grupo del cliente", False),
    ("Asignar tecnico", "Sin asignar (opcional)", False),
]
for i, (label, placeholder, req) in enumerate(fields):
    y = 2.0 + i * 0.75
    req_mark = " *" if req else ""
    txt(s, 0.8, y, 6.4, 0.25, label + req_mark, 12, DARK, True)
    field = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y+0.25), Inches(6.4), Inches(0.35))
    field.fill.solid(); field.fill.fore_color.rgb = RGBColor(0xf0, 0xf2, 0xf8); field.line.color.rgb = RGBColor(0xd5, 0xdb, 0xe6); field.line.width = Pt(0.5)
    txt(s, 0.95, y+0.28, 6, 0.3, placeholder, 10, MUTED)

# Buttons
btn_save = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(6.4), Inches(2.1), Inches(0.4))
btn_save.fill.solid(); btn_save.fill.fore_color.rgb = ACCENT; btn_save.line.fill.background()
txt(s, 5.3, 6.42, 2.1, 0.35, "Guardar reporte", 12, WHITE, True, PP_ALIGN.CENTER)
btn_cancel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3), Inches(6.4), Inches(2), Inches(0.4))
btn_cancel.fill.solid(); btn_cancel.fill.fore_color.rgb = RGBColor(0xf0, 0xf2, 0xf8); btn_cancel.line.fill.background()
txt(s, 3, 6.42, 2, 0.35, "Cancelar", 12, DARK, False, PP_ALIGN.CENTER)

# Right side explanation
txt(s, 8, 1.5, 5, 0.4, "Campos del formulario:", 16, TEXT, True)
items = [
    "Cliente: nombre del negocio o persona",
    "Equipo: marca, modelo, serie del equipo",
    "Descripcion: texto libre del problema",
    "Prioridad: baja, normal, alta, urgente",
    "Grupo WhatsApp: se envia automatico al grupo",
    "Tecnico: se puede asignar al crear",
    "",
    "Al guardar se crea en estado 'Abierto'",
    "Se puede enviar por WhatsApp despues",
    "El reporte queda visible para el tecnico"
]
bullets(s, 8, 2, 5, 4.5, items, 12, TEXT2)

# ===== SLIDE 8: DESKTOP - CLIENTES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Desktop: Clientes y Equipos", "Gestion completa de la base de datos")

# Clients panel
txt(s, 0.4, 1.2, 6, 0.3, "Gestion de clientes", 16, ACCENT, True)
# Search bar
sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.6), Inches(4), Inches(0.35))
sb.fill.solid(); sb.fill.fore_color.rgb = SURFACE; sb.line.fill.background()
txt(s, 0.55, 1.63, 3.7, 0.3, "Buscar cliente...", 10, MUTED)
btn_nc = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.6), Inches(1.6), Inches(1.8), Inches(0.35))
btn_nc.fill.solid(); btn_nc.fill.fore_color.rgb = ACCENT; btn_nc.line.fill.background()
txt(s, 4.6, 1.63, 1.8, 0.3, "+ Nuevo cliente", 10, WHITE, True, PP_ALIGN.CENTER)

clients = [
    ("tacos El Güero", "555-1234", "Grupo WA: tacos-gerencia"),
    ("Cafe Azul", "555-5678", "Grupo WA: cafe-azul-ops"),
    ("Farmacia Lopez", "555-9012", "Grupo WA: farmacia-lopez"),
    ("AutoZone Centro", "555-3456", "Sin grupo WA"),
]
for i, (name, tel, grupo) in enumerate(clients):
    y = 2.1 + i*0.65
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(6), Inches(0.55))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    txt(s, 0.55, y+0.05, 3, 0.2, name, 11, TEXT, True)
    txt(s, 0.55, y+0.25, 3, 0.18, tel, 9, MUTED)
    txt(s, 3.5, y+0.1, 2.8, 0.2, grupo, 9, GREEN if "WA" in grupo else MUTED)

# Equipment panel
txt(s, 6.8, 1.2, 6, 0.3, "Equipos por cliente", 16, ACCENT, True)
eq_sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.6), Inches(4), Inches(0.35))
eq_sb.fill.solid(); eq_sb.fill.fore_color.rgb = SURFACE; eq_sb.line.fill.background()
txt(s, 6.95, 1.63, 3.7, 0.3, "Buscar equipo...", 10, MUTED)
btn_ne = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11), Inches(1.6), Inches(1.8), Inches(0.35))
btn_ne.fill.solid(); btn_ne.fill.fore_color.rgb = ACCENT; btn_ne.line.fill.background()
txt(s, 11, 1.63, 1.8, 0.3, "+ Nuevo equipo", 10, WHITE, True, PP_ALIGN.CENTER)

equipos = [
    ("Freidora #3", "tacos El Güero", "8 reportes"),
    ("Computo #1", "Cafe Azul", "5 reportes"),
    ("Refrigerador #2", "Farmacia Lopez", "4 reportes"),
    ("Impresora Termica", "AutoZone Centro", "2 reportes"),
]
for i, (name, client, reps) in enumerate(equipos):
    y = 2.1 + i*0.65
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(y), Inches(6), Inches(0.55))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    txt(s, 6.95, y+0.05, 3, 0.2, name, 11, TEXT, True)
    txt(s, 6.95, y+0.25, 3, 0.18, client, 9, MUTED)
    pill(s, 10.5, y+0.15, reps, TEXT, SURFACE2, 1.2)

txt(s, 0.4, 5, 12.5, 0.5, "El gerente registra clientes y equipos. El tecnico los usa al crear reportes.", 12, MUTED, False, PP_ALIGN.CENTER)
txt(s, 0.4, 5.4, 12.5, 0.5, "Cada cliente puede tener un grupo de WhatsApp vinculado para envio automatico de reportes.", 12, MUTED, False, PP_ALIGN.CENTER)

# ===== SLIDE 9: DESKTOP - TECNICOS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Desktop: Tecnicos", "Administrar el equipo de trabajo")

txt(s, 0.4, 1.2, 6, 0.3, "Lista de tecnicos", 16, ACCENT, True)
btn_import = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8), Inches(1.2), Inches(2.5), Inches(0.35))
btn_import.fill.solid(); btn_import.fill.fore_color.rgb = SURFACE2; btn_import.line.fill.background()
txt(s, 8, 1.23, 2.5, 0.3, "Importar desde archivo", 10, TEXT2, False, PP_ALIGN.CENTER)
btn_nt = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.6), Inches(1.2), Inches(2.1), Inches(0.35))
btn_nt.fill.solid(); btn_nt.fill.fore_color.rgb = ACCENT; btn_nt.line.fill.background()
txt(s, 10.6, 1.23, 2.1, 0.3, "+ Nuevo tecnico", 10, WHITE, True, PP_ALIGN.CENTER)

techs = [
    ("Juan Perez", "juan", "555-1001", "11 reportes", True),
    ("Pedro Lopez", "pedro", "555-1002", "10 reportes", True),
    ("Maria Garcia", "maria", "555-1003", "8 reportes", True),
    ("Carlos Ruiz", "carlos", "555-1004", "3 reportes", False),
]
for i, (name, user, tel, reps, active) in enumerate(techs):
    y = 1.7 + i*0.7
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(y), Inches(12.4), Inches(0.6))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    if not active:
        card.fill.fore_color.rgb = RGBColor(0x18, 0x18, 0x20)
    txt(s, 0.6, y+0.05, 2.5, 0.2, name, 12, TEXT if active else MUTED, True)
    txt(s, 0.6, y+0.28, 2.5, 0.18, "Usuario: " + user + " | " + tel, 9, MUTED)
    pill(s, 3.5, y+0.15, reps, TEXT, SURFACE2, 1.2)
    pill(s, 5, y+0.15, "Activo" if active else "Inactivo", GREEN if active else RED, GREEN_BG if active else RED_BG)
    pill(s, 6.5, y+0.15, "Tecnico", ACCENT, BLUE_BG)

txt(s, 0.4, 4.6, 12.5, 0.4, "Cada tecnico tiene: nombre, usuario, contrasena, rol (tecnico/gerente)", 12, MUTED, False, PP_ALIGN.CENTER)
txt(s, 0.4, 5, 12.5, 0.4, "Se pueden importar masivamente desde archivo  |  Los tecnicos inactivos no pueden loguearse", 12, MUTED, False, PP_ALIGN.CENTER)

# ===== SLIDE 10: DESKTOP - MENSAJES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Desktop: Mensajes WhatsApp", "Chat integrado estilo WhatsApp Web")

# Side panel
side = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.3), Inches(1.3), Inches(3.5), Inches(5.7))
side.fill.solid(); side.fill.fore_color.rgb = WHITE; side.line.fill.background()
# Side header
side_head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.3), Inches(3.5), Inches(0.45))
side_head.fill.solid(); side_head.fill.fore_color.rgb = RGBColor(0x00, 0x80, 0x69); side_head.line.fill.background()
txt(s, 0.5, 1.35, 2, 0.35, "Mensajes", 14, WHITE, True)
# Search
ss = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.85), Inches(3.1), Inches(0.3))
ss.fill.solid(); ss.fill.fore_color.rgb = RGBColor(0xf0, 0xf2, 0xf5); ss.line.fill.background()
txt(s, 0.65, 1.88, 2.8, 0.25, "Buscar o empezar un chat...", 9, MUTED)

chats = [
    ("tacos El Güero", "Grupo", "Juan: Freidora arreglada", "2:30 pm"),
    ("Cafe Azul", "", "Pedido de bombillo LED", "1:15 pm"),
    ("Farmacia Lopez", "Grupo", "Maria: Esperando repuesto", "11:00 am"),
    ("+52 555-1234", "", "Hola, cuando vienen?", "Ayer"),
]
for i, (name, tipo, preview, time) in enumerate(chats):
    y = 2.25 + i*0.7
    # Avatar
    ava = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(y+0.05), Inches(0.45), Inches(0.45))
    ava.fill.solid(); ava.fill.fore_color.rgb = ACCENT if tipo == "Grupo" else GREEN; ava.line.fill.background()
    txt(s, 0.5, y+0.1, 0.45, 0.35, name[0], 12, WHITE, True, PP_ALIGN.CENTER)
    txt(s, 1.05, y+0.03, 2, 0.2, name + (" [grupo]" if tipo else ""), 10, DARK, True)
    txt(s, 1.05, y+0.25, 2.3, 0.18, preview, 8, MUTED)
    txt(s, 3, y+0.05, 0.6, 0.18, time, 7, MUTED, False, PP_ALIGN.RIGHT)

# Chat area
chat_area = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.3), Inches(9), Inches(5.7))
chat_area.fill.solid(); chat_area.fill.fore_color.rgb = RGBColor(0xef, 0xea, 0xe2); chat_area.line.fill.background()
# Chat header
chat_head = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(1.3), Inches(9), Inches(0.45))
chat_head.fill.solid(); chat_head.fill.fore_color.rgb = RGBColor(0xf0, 0xf2, 0xf5); chat_head.line.fill.background()
txt(s, 4.2, 1.35, 4, 0.35, "tacos El Güero (grupo)  |  4 participantes", 11, DARK, True)

# Messages
msgs = [
    (True, "Hola, la freidora #3 no enciende", "10:30 am"),
    (False, "Juan: Ya voy a revisar", "10:35 am"),
    (True, "Es que ayer funcionaba bien", "10:37 am"),
    (False, "Juan: Encontre el problema, es el fusible", "11:00 am"),
    (False, "Juan: [Imagen del fusible]", "11:01 am"),
    (True, "Gracias! Cuando queda?", "11:05 am"),
]
for i, (out, text, time) in enumerate(msgs):
    y = 1.9 + i*0.5
    if out:
        # Outgoing (right aligned)
        bubble = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.5), Inches(y), Inches(4.5), Inches(0.4))
        bubble.fill.solid(); bubble.fill.fore_color.rgb = RGBColor(0xd9, 0xfd, 0xd3); bubble.line.fill.background()
        txt(s, 7.6, y+0.02, 4.3, 0.18, text, 9, DARK)
        txt(s, 11, y+0.22, 1, 0.15, time + " ✓✓", 7, MUTED, False, PP_ALIGN.RIGHT)
    else:
        # Incoming (left aligned)
        bubble = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(y), Inches(4.5), Inches(0.4))
        bubble.fill.solid(); bubble.fill.fore_color.rgb = WHITE; bubble.line.fill.background()
        txt(s, 4.3, y+0.02, 4.3, 0.18, text, 9, DARK)
        txt(s, 7.3, y+0.22, 1.2, 0.15, time, 7, MUTED, False, PP_ALIGN.RIGHT)

# Reply bar
reply_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(6.55), Inches(9), Inches(0.45))
reply_bar.fill.solid(); reply_bar.fill.fore_color.rgb = RGBColor(0xf0, 0xf2, 0xf5); reply_bar.line.fill.background()
txt(s, 4.2, 6.58, 6, 0.35, "Escribe un mensaje", 11, MUTED)
# Send button
send_btn = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12.3), Inches(6.6), Inches(0.35), Inches(0.35))
send_btn.fill.solid(); send_btn.fill.fore_color.rgb = RGBColor(0x00, 0xa8, 0x84); send_btn.line.fill.background()
txt(s, 12.3, 6.63, 0.35, 0.3, ">", 12, WHITE, True, PP_ALIGN.CENTER)

txt(s, 0.3, 7.1, 12.7, 0.25, "Funciones: enviar texto, fotos, videos, stickers, notas de voz, emojis, menciones @, responder, crear reportes desde chats", 10, MUTED, False, PP_ALIGN.CENTER)

# ===== SLIDE 11: MOBILE - LOGIN =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "App Movil: Login y Biometria", "Primera vez y acceso rapido")

# Login screen mockup
inner = phone_mockup(s, 0.5, 1.3, 4, 5.8)
# Logo
logo_circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.8), Inches(1.6), Inches(1.3), Inches(1.3))
logo_circle.fill.solid(); logo_circle.fill.fore_color.rgb = SURFACE; logo_circle.line.fill.background()
txt(s, 1.8, 1.8, 1.3, 0.8, "NA", 28, ACCENT, True, PP_ALIGN.CENTER)
txt(s, 0.7, 3.05, 3.6, 0.3, "Bienvenido, tecnico", 12, TEXT, True, PP_ALIGN.CENTER)

# Form fields
txt(s, 0.9, 3.5, 3.2, 0.2, "Usuario", 9, MUTED, True)
uf = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(3.7), Inches(3.2), Inches(0.4))
uf.fill.solid(); uf.fill.fore_color.rgb = RGBColor(0x0f, 0x15, 0x26); uf.line.fill.background()
txt(s, 1.05, 3.73, 2.9, 0.35, "Tu usuario", 10, MUTED)

txt(s, 0.9, 4.2, 3.2, 0.2, "Contrasena", 9, MUTED, True)
pf = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.4), Inches(3.2), Inches(0.4))
pf.fill.solid(); pf.fill.fore_color.rgb = RGBColor(0x0f, 0x15, 0x26); pf.line.fill.background()
txt(s, 1.05, 4.43, 2.9, 0.35, "Tu contrasena", 10, MUTED)

# Login button
login_btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(4.95), Inches(3.2), Inches(0.45))
login_btn.fill.solid(); login_btn.fill.fore_color.rgb = ACCENT; login_btn.line.fill.background()
txt(s, 0.9, 4.98, 3.2, 0.4, "Entrar", 13, WHITE, True, PP_ALIGN.CENTER)

# Biometric button
bio_btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.9), Inches(5.55), Inches(3.2), Inches(0.45))
bio_btn.fill.solid(); bio_btn.fill.fore_color.rgb = SURFACE; bio_btn.line.color.rgb = BORDER; bio_btn.line.width = Pt(0.5)
txt(s, 0.9, 5.58, 3.2, 0.4, "Entrar con biometria", 11, TEXT, False, PP_ALIGN.CENTER)
txt(s, 1.3, 5.95, 2.4, 0.2, "Huella o rostro", 8, MUTED, False, PP_ALIGN.CENTER)

# Right side explanation
txt(s, 5.2, 1.5, 7.5, 0.4, "Proceso de login:", 18, TEXT, True)
steps = [
    ("1. Escribe tu usuario", "El gerente te asigna usuario y contrasena"),
    ("2. Guarda tu biometria", "La primera vez toca 'biometria' y registra tu huella"),
    ("3. En las proximas veces", "Toca el boton de biometria y listo, sin escribir nada"),
    ("4. Si falla 3 veces", "Reinicia la app e intenta de nuevo"),
]
for i, (step, desc) in enumerate(steps):
    y = 2.1 + i*0.85
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.2), Inches(y), Inches(7.5), Inches(0.7))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.4), Inches(y+0.12), Inches(0.45), Inches(0.45))
    num.fill.solid(); num.fill.fore_color.rgb = ACCENT; num.line.fill.background()
    txt(s, 5.4, y+0.15, 0.45, 0.4, str(i+1), 14, WHITE, True, PP_ALIGN.CENTER)
    txt(s, 6, y+0.08, 6.5, 0.25, step, 12, TEXT, True)
    txt(s, 6, y+0.35, 6.5, 0.25, desc, 10, MUTED)

# ===== SLIDE 12: MOBILE - REPORTE DETALLE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "App Movil: Detalle del Reporte", "Toda la informacion en una pantalla")

# Phone with detail view
inner = phone_mockup(s, 0.3, 1.3, 4.3, 6)
# Topbar
tb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.36), Inches(1.35), Inches(4.18), Inches(0.35))
tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0x0a, 0x0e, 0x1a); tb.line.fill.background()
txt(s, 0.5, 1.37, 2, 0.3, "NexAlert", 10, ACCENT, True)
txt(s, 0.5, 1.7, 2, 0.25, "tacos El Güero", 13, TEXT, True)
txt(s, 0.5, 1.95, 2, 0.2, "Freidora #3", 9, MUTED)
pill(s, 2.8, 1.72, "Abierto", WHITE, RED, 1)
pill(s, 2.8, 1.95, "Urgente", WHITE, RED, 1)

# Detail rows
detail_rows = [
    ("FECHA DEL REPORTE", "20 ago 2026"),
    ("PROBLEMA", "No enciende, se escucha ruido raro cuando se conecta a la corriente"),
    ("UBICACION", "Abrir en Maps (19.4326, -99.1332)"),
    ("EVIDENCIA (2)", "[foto1.jpg] [foto2.jpg]"),
]
y = 2.25
for label, value in detail_rows:
    txt(s, 0.5, y, 3.9, 0.12, label, 7, MUTED, True)
    txt(s, 0.5, y+0.13, 3.9, 0.25, value, 8, TEXT2)
    y += 0.45

# Photo buttons
cam_btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(1.85), Inches(0.35))
cam_btn.fill.solid(); cam_btn.fill.fore_color.rgb = ACCENT; cam_btn.line.fill.background()
txt(s, 0.5, y+0.02, 1.85, 0.3, "Camara", 9, WHITE, True, PP_ALIGN.CENTER)
gal_btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(y), Inches(1.85), Inches(0.35))
gal_btn.fill.solid(); gal_btn.fill.fore_color.rgb = SURFACE; gal_btn.line.fill.background()
txt(s, 2.5, y+0.02, 1.85, 0.3, "Galeria", 9, TEXT2, False, PP_ALIGN.CENTER)
y += 0.45

# Action buttons
actions = ["Compartir", "Ubicacion", "Historial"]
for i, act in enumerate(actions):
    abtn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i*1.35), Inches(y), Inches(1.2), Inches(0.3))
    abtn.fill.solid(); abtn.fill.fore_color.rgb = SURFACE; abtn.line.fill.background()
    txt(s, 0.5 + i*1.35, y+0.02, 1.2, 0.25, act, 8, TEXT2, False, PP_ALIGN.CENTER)
y += 0.4

# Status buttons
estados = ["Abierto", "En proceso", "Espera rep.", "Espera cli.", "Resuelto"]
est_colors = [RED, AMBER, MUTED, MUTED, GREEN]
for i, (est, ec) in enumerate(zip(estados, est_colors)):
    ebtn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5 + i*0.75), Inches(y), Inches(0.7), Inches(0.3))
    if est == "Abierto":
        ebtn.fill.solid(); ebtn.fill.fore_color.rgb = ACCENT; ebtn.line.fill.background()
    else:
        ebtn.fill.solid(); ebtn.fill.fore_color.rgb = SURFACE; ebtn.line.fill.background()
    txt(s, 0.5 + i*0.75, y+0.02, 0.7, 0.25, est, 6, WHITE if est=="Abierto" else MUTED, est=="Abierto", PP_ALIGN.CENTER)
y += 0.45

# Notes
txt(s, 0.5, y, 3.9, 0.15, "COMENTARIOS (1)", 7, MUTED, True)
y += 0.18
nota = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(3.9), Inches(0.45))
nota.fill.solid(); nota.fill.fore_color.rgb = SURFACE; nota.line.fill.background()
txt(s, 0.6, y+0.02, 3.7, 0.2, "Juan: Ya estoy en camino", 8, TEXT)
txt(s, 0.6, y+0.23, 3.7, 0.15, "Juan Perez · 20 ago, 10:35", 7, MUTED)

# Right side
txt(s, 5, 1.5, 8, 0.4, "Pantalla de detalle:", 18, TEXT, True)
bullets(s, 5, 2, 8, 5, [
    "Cabecera: cliente, equipo, badges de estado y prioridad",
    "Fecha del reporte",
    "Descripcion del problema (texto libre)",
    "Solucion (se llena al resolver)",
    "Ubicacion GPS con link a Google Maps",
    "Evidencia fotografica (galeria de fotos)",
    "Botones: Camara y Galeria para subir fotos",
    "Compartir: envia el reporte formateado por WhatsApp",
    "Registrar ubicacion: captura GPS actual",
    "Historial: eventos del reporte (cambios, notas, fotos)",
    "",
    "Botones de estado (cambiar con un toque):",
    "  Abierto → En proceso → Resuelto",
    "  Espera repuesto / Espera cliente (requiere nota)",
    "",
    "Cuando se resuelve aparece:",
    "  Firma del cliente (canvas touch)",
    "  Campo de solucion",
    "",
    "Comentarios: notas del tecnico con fecha/hora",
    "Boton para enviar nuevo comentario"
], 11, TEXT2, 3)

# ===== SLIDE 13: MOBILE - CREAR REPORTE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "App Movil: Crear Reporte", "El tecnico crea reportes desde el campo")

# Phone with form
inner = phone_mockup(s, 0.3, 1.3, 4.3, 6)
tb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.36), Inches(1.35), Inches(4.18), Inches(0.35))
tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0x0a, 0x0e, 0x1a); tb.line.fill.background()
txt(s, 0.5, 1.37, 2, 0.3, "< Volver", 10, ACCENT, True)
txt(s, 0.5, 1.72, 3.9, 0.3, "Nuevo reporte", 14, TEXT, True)

form_fields = [
    ("Cliente *", "tacos El Güero"),
    ("Equipo", "Freidora #3"),
    ("Descripcion *", "No enciende, ruido raro"),
    ("Prioridad", "Normal"),
    ("Asignar a", "Sin asignar"),
]
y = 2.15
for label, val in form_fields:
    txt(s, 0.5, y, 3.9, 0.15, label, 8, MUTED, True)
    field = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y+0.17), Inches(3.9), Inches(0.35))
    field.fill.solid(); field.fill.fore_color.rgb = RGBColor(0x0f, 0x15, 0x26); field.line.fill.background()
    txt(s, 0.65, y+0.2, 3.6, 0.3, val, 10, TEXT)
    y += 0.58

save_btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(3.9), Inches(0.45))
save_btn.fill.solid(); save_btn.fill.fore_color.rgb = ACCENT; save_btn.line.fill.background()
txt(s, 0.5, y+0.03, 3.9, 0.4, "Guardar reporte", 12, WHITE, True, PP_ALIGN.CENTER)

txt(s, 5, 1.5, 8, 0.4, "Crear reporte desde el celular:", 18, TEXT, True)
bullets(s, 5, 2, 8, 5.5, [
    "Cliente: escribir nombre o buscar existente",
    "Equipo: marca, modelo, serie del equipo",
    "Descripcion: texto libre del problema (requerido)",
    "Prioridad: Baja / Normal / Alta / Urgente",
    "Asignar tecnico: opcional, puede dejarse vacio",
    "",
    "Al guardar el reporte:",
    "  → Se crea en estado 'Abierto'",
    "  → Se registra la fecha automaticamente",
    "  → Si hay internet, se envia al servidor",
    "  → Si NO hay internet, se guarda local (IndexedDB)",
    "  → Se sincroniza cuando vuelva la conexion",
    "",
    "Despues de crear puedes:",
    "  → Abrir el detalle para agregar fotos",
    "  → Cambiar el estado",
    "  → Agregar comentarios",
    "  → Compartir por WhatsApp",
    "  → Registrar ubicacion GPS"
], 11, TEXT2, 3)

# ===== SLIDE 14: ESTADOS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Estados de un reporte", "Flujo completo de vida de un reporte")

estados_info = [
    ("Abierto", "Reporte nuevo, sin atender", RED, "Se crea automaticamente al crear el reporte"),
    ("En Proceso", "Tecnico esta trabajando", AMBER, "El tecnico avanza que ya esta en el tema"),
    ("Espera Repuesto", "Necesita una pieza", MUTED, "Requiere nota: que pieza se necesita"),
    ("Espera Cliente", "Info del cliente", MUTED, "Requiere nota: que informacion se necesita"),
    ("Resuelto", "Equipo reparado", GREEN, "Aparece firma del cliente + campo solucion"),
]
for i, (nome, desc, cor, detail) in enumerate(estados_info):
    left = 0.3 + i * 2.55
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.4), Inches(2.35), Inches(2.5))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    top_line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(1.4), Inches(2.35), Inches(0.05))
    top_line.fill.solid(); top_line.fill.fore_color.rgb = cor; top_line.line.fill.background()
    txt(s, left+0.1, 1.6, 2.15, 0.3, nome, 14, cor, True, PP_ALIGN.CENTER)
    txt(s, left+0.1, 2.0, 2.15, 0.35, desc, 10, TEXT2, False, PP_ALIGN.CENTER)
    txt(s, left+0.1, 2.5, 2.15, 0.8, detail, 9, MUTED, False, PP_ALIGN.CENTER)
    if i < 4:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left+2.4), Inches(2.3), Inches(0.2), Inches(0.2))
        arr.fill.solid(); arr.fill.fore_color.rgb = MUTED; arr.line.fill.background()

# Prioridades
txt(s, 0.5, 4.2, 12, 0.4, "Prioridades:", 16, TEXT, True)
prios = [
    ("Urgente", "Requiere atencion inmediata", WHITE, RED),
    ("Alta", "Importante, resolver pronto", WHITE, ORANGE),
    ("Normal", "Flujo regular", WHITE, ACCENT),
    ("Baja", "Puede esperar", WHITE, MUTED),
]
for i, (nome, desc, tc, bc) in enumerate(prios):
    left = 0.5 + i * 3.1
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(4.7), Inches(2.8), Inches(0.8))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    pill(s, left+0.1, 4.8, nome, tc, bc, 1.3)
    txt(s, left+0.1, 5.15, 2.6, 0.25, desc, 10, MUTED)

txt(s, 0.5, 5.8, 12, 0.4, "Regla de archivado: Los reportes en estado 'Resuelto' se archivan automaticamente despues de 24 horas", 12, ACCENT, False, PP_ALIGN.CENTER)
txt(s, 0.5, 6.2, 12, 0.4, "Los reportes archivados se pueden ver en la seccion 'Archivados' del menu", 12, MUTED, False, PP_ALIGN.CENTER)

# ===== SLIDE 15: WHATSAPP =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Integracion WhatsApp", "Conexion directa con clientes")

# Desktop QR flow
txt(s, 0.5, 1.2, 6, 0.3, "Conexion en el Desktop:", 14, ACCENT, True)
steps_wa = [
    ("1. Click en 'Conectar WhatsApp'", "Se genera un codigo QR"),
    ("2. Escanea con tu telefono", "WhatsApp > Ajustes > Dispositivos vinculados"),
    ("3. Listo, ya estas conectado", "El indicador se pone verde"),
    ("4. Selecciona un grupo", "Vincula el grupo al cliente"),
]
for i, (step, desc) in enumerate(steps_wa):
    y = 1.6 + i*0.55
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(5.8), Inches(0.45))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    num = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), Inches(y+0.07), Inches(0.3), Inches(0.3))
    num.fill.solid(); num.fill.fore_color.rgb = GREEN; num.line.fill.background()
    txt(s, 0.6, y+0.08, 0.3, 0.28, str(i+1), 9, WHITE, True, PP_ALIGN.CENTER)
    txt(s, 1, y+0.05, 3, 0.2, step, 10, TEXT, True)
    txt(s, 1, y+0.23, 5, 0.18, desc, 8, MUTED)

# Chat features
txt(s, 6.8, 1.2, 6, 0.3, "Funciones del chat:", 14, ACCENT, True)
chat_feats = [
    "Enviar y recibir mensajes de texto",
    "Enviar fotos, videos, stickers, documentos",
    "Notas de voz (grabar y enviar)",
    "Emojis completos (panel desplegable)",
    "Menciones @ en grupos",
    "Responder a mensajes especificos",
    "Crear reportes desde mensajes seleccionados",
    "Seleccionar multiples mensajes",
    "Borrar mensajes",
    "Busqueda de conversaciones",
    "Avatars con iniciales o fotos",
    "Separadores de fecha (Hoy, Ayer, etc)",
]
bullets(s, 7, 1.6, 6, 4.5, chat_feats, 11, TEXT2, 3)

# Share report format
txt(s, 0.5, 4, 12, 0.3, "Formato del reporte enviado por WhatsApp:", 14, ACCENT, True)
report_format = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(4.4), Inches(6), Inches(2.5))
report_format.fill.solid(); report_format.fill.fore_color.rgb = RGBColor(0xdc, 0xf8, 0xc6); report_format.line.fill.background()
txt(s, 0.7, 4.5, 5.6, 2.3, "REPORTE DE FALLA\n\nPrioridad: Urgente\nEmpresa: tacos El Güero\nEquipo: Freidora #3\nFecha: 20 ago 2026\nProblema:\nNo enciende, se escucha\nruido raro al conectar\n\n— NexAlert", 11, DARK)

txt(s, 7, 4.4, 6, 0.3, "Desde el movil:", 13, TEXT, True)
bullets(s, 7, 4.8, 6, 2, [
    "Toca 'Compartir' en el detalle del reporte",
    "Se abre el selector de WhatsApp",
    "Elige el contacto o grupo",
    "Se envia el mensaje formateado"
], 11, TEXT2, 4)

# ===== SLIDE 16: OFFLINE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Modo Offline", "La app funciona sin conexion a internet")

for i, (num, title_txt, desc, col) in enumerate([
    ("1", "Sin internet", "Crea reportes,\ncambia estados,\nagrega fotos", RED),
    ("2", "Guardado local", "IndexedDB almacena\ntodos los cambios\nen el celular", AMBER),
    ("3", "Sincronizacion", "Cuando vuelve\nel internet, se\nenvian los cambios", GREEN),
    ("4", "Cola de fotos", "Las fotos se\ncomprimen y se\nsuben despues", ACCENT),
]):
    left = 0.5 + i * 3.2
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.3), Inches(2.9), Inches(3.2))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.9), Inches(1.5), Inches(1.1), Inches(1.1))
    circ.fill.solid(); circ.fill.fore_color.rgb = col; circ.line.fill.background()
    txt(s, left + 0.9, 1.6, 1.1, 0.9, num, 28, WHITE, True, PP_ALIGN.CENTER)
    txt(s, left + 0.2, 2.75, 2.5, 0.35, title_txt, 15, col, True, PP_ALIGN.CENTER)
    txt(s, left + 0.2, 3.2, 2.5, 1.2, desc, 11, TEXT2, False, PP_ALIGN.CENTER)
    if i < 3:
        arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left + 3.0), Inches(2.6), Inches(0.25), Inches(0.25))
        arr.fill.solid(); arr.fill.fore_color.rgb = col; arr.line.fill.background()

# Queue details
txt(s, 0.5, 4.8, 12, 0.3, "Cola de cambios pendientes:", 14, ACCENT, True)
queue_items = [
    ("Cambios de estado", "Se guardan local y se envian cuando haya internet"),
    ("Comentarios / Notas", "Se almacenan en IndexedDB y se sincronizan"),
    ("Fotos", "Se comprimen a 1280px JPEG 72% y se suben en segundo plano"),
    ("Ubicacion GPS", "Se captura y se envia con el reporte"),
    ("Reportes nuevos", "Se crean local y se sincronizan al servidor"),
    ("Badge de pendientes", "Muestra cuantos cambios esperan en la barra superior"),
]
for i, (item, desc) in enumerate(queue_items):
    col = i % 3
    row = i // 3
    left = 0.5 + col * 4.2
    top = 5.2 + row * 0.9
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.9), Inches(0.75))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    txt(s, left+0.15, top+0.05, 3.6, 0.2, item, 11, TEXT, True)
    txt(s, left+0.15, top+0.28, 3.6, 0.4, desc, 9, MUTED)

# ===== SLIDE 17: MONITOREO =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Monitoreo (Gerente en Movil)", "Supervisar tecnicos desde el celular")

# Phone with monitoreo
inner = phone_mockup(s, 0.3, 1.3, 4.3, 6)
tb = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.36), Inches(1.35), Inches(4.18), Inches(0.35))
tb.fill.solid(); tb.fill.fore_color.rgb = RGBColor(0x0a, 0x0e, 0x1a); tb.line.fill.background()
txt(s, 0.5, 1.37, 2, 0.3, "< Volver", 10, ACCENT, True)
txt(s, 0.5, 1.72, 3.9, 0.3, "Monitoreo", 14, TEXT, True)

# Stats
for i, (n, l, c) in enumerate([(47,"Total",ACCENT),(7,"Abiertos",RED),(6,"Proceso",AMBER),(24,"Resueltos",GREEN)]):
    stat_card(s, 0.5 + i*0.95, 2.1, n, l, c)

# Tech filters
txt(s, 0.5, 3.05, 3.9, 0.15, "Filtrar por tecnico:", 8, MUTED, True)
techs_m = [("Juan Perez", 3, 5), ("Pedro Lopez", 2, 8), ("Maria Garcia", 1, 4), ("Sin asignar", 5, 0)]
for i, (name, abiertos, resueltos) in enumerate(techs_m):
    y = 3.25 + i*0.4
    row = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(y), Inches(3.9), Inches(0.35))
    row.fill.solid(); row.fill.fore_color.rgb = SURFACE; row.line.fill.background()
    txt(s, 0.65, y+0.03, 2, 0.15, name, 8, TEXT, True)
    txt(s, 2.8, y+0.05, 1.5, 0.15, str(abiertos) + " abiertos · " + str(resueltos) + " resueltos", 7, MUTED)

txt(s, 5, 1.5, 8, 0.4, "Acceso al monitoreo:", 18, TEXT, True)
bullets(s, 5, 2, 8, 4.5, [
    "Solo el gerente ve la opcion 'Monitoreo' en el menu",
    "Se accede desde el hamburger menu (3 puntos)",
    "",
    "Datos que muestra:",
    "  Estadisticas globales (total, abiertos, proceso, resueltos)",
    "  Lista de tecnicos con sus reportes",
    "  Filtro por tecnico para ver solo los de uno",
    "  Filtro 'Sin asignar' para ver reportes sin tecnico",
    "",
    "Acciones del gerente desde el movil:",
    "  Ver detalle de cualquier reporte",
    "  Cambiar estado de reportes",
    "  Asignar/reasignar tecnicos",
    "  Desactivar tecnicos",
    "  Ver historial de eventos",
    "",
    "Notificaciones push:",
    "  El gerente recibe push cuando se crea un reporte nuevo",
    "  Y cuando se asigna un reporte a un tecnico"
], 11, TEXT2, 3)

# ===== SLIDE 18: NOTIFICACIONES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Notificaciones Push", "Alertas en tiempo real via Firebase")

# Phone with notification
inner = phone_mockup(s, 0.5, 1.5, 3.5, 4.5)
# Notification popup
notif = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.7), Inches(3.1), Inches(1))
notif.fill.solid(); notif.fill.fore_color.rgb = WHITE; notif.line.fill.background()
txt(s, 0.85, 1.78, 2.8, 0.15, "NexAlert", 8, ACCENT, True)
txt(s, 0.85, 1.95, 2.8, 0.15, "Reporte asignado", 10, DARK, True)
txt(s, 0.85, 2.15, 2.8, 0.25, "tacos El Güero - Freidora #3\nNo enciende, se escucha ruido", 8, MUTED)

txt(s, 5, 1.5, 8, 0.4, "Cuando se envian notificaciones:", 18, TEXT, True)
push_events = [
    ("Nuevo reporte creado", "El gerente recibe: 'Nuevo reporte: cliente - equipo'"),
    ("Reporte asignado", "El tecnico recibe: 'Reporte asignado: cliente - equipo'"),
    ("Cambio de estado", "Se notifica al gerente del cambio"),
    ("Recordatorio", "Reportes sin actualizacion se record cada 5 min"),
]
for i, (event, desc) in enumerate(push_events):
    y = 2.1 + i*0.7
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(y), Inches(8), Inches(0.6))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(y+0.15), Inches(0.3), Inches(0.3))
    dot.fill.solid(); dot.fill.fore_color.rgb = GREEN; dot.line.fill.background()
    txt(s, 5.7, y+0.08, 7, 0.2, event, 12, TEXT, True)
    txt(s, 5.7, y+0.3, 7, 0.2, desc, 10, MUTED)

txt(s, 5, 5, 8, 0.3, "Configuracion:", 14, TEXT, True)
bullets(s, 5, 5.4, 8, 1.5, [
    "El registro de push se hace automaticamente al hacer login",
    "El servidor almacena los tokens Firebase por tecnico",
    "Se pueden probar notificaciones desde el panel de ajustes",
    "Configuracion del celular: Ajustes > NexAlert > Notificaciones > Activar"
], 11, TEXT2, 4)

# ===== SLIDE 19: PROBLEMAS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, BG); topbar(s)
title(s, "Problemas comunes y soluciones", "Lo mas frecuente que puede pasar")

problemas = [
    ("No puedo entrar", "Verifica usuario/contrasena. Si falla 3 veces, reinicia la app.", RED),
    ("No me sale biometria", "Ve al menu > Tema y busca la opcion. Asegurate de tener huella en el celular.", AMBER),
    ("No se envia el reporte", "Revisa tu conexion. El reporte se guardo local y se enviara despues.", AMBER),
    ("No recibo notificaciones", "Ve a Ajustes > NexAlert > Notificaciones y activalas.", RED),
    ("La app esta lenta", "Cierra otras apps. Si persiste, reinicia el celular.", MUTED),
    ("WhatsApp no conecta", "Reinicia sesion desde el Desktop. Escanea el QR de nuevo.", AMBER),
    ("Las fotos no suben", "Estan en cola. Se subiran cuando haya internet.", MUTED),
    ("No veo reportes", "Toca 'Actualizar' en el menu. O verifica tu conexion.", RED),
]
for i, (prob, sol, col) in enumerate(problemas):
    col_i = i % 2
    row_i = i // 2
    left = 0.3 + col_i * 6.5
    top = 1.2 + row_i * 1.5
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(6.2), Inches(1.3))
    card.fill.solid(); card.fill.fore_color.rgb = SURFACE; card.line.fill.background()
    txt(s, left+0.2, top+0.1, 5.8, 0.25, prob, 13, col, True)
    txt(s, left+0.2, top+0.4, 5.8, 0.7, sol, 10, TEXT2)

# ===== SLIDE 20: CIERRE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, DARK); topbar(s)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5))
txt(s, 1, 3, 11.3, 0.8, "Listo para empezar!", 44, WHITE, True, PP_ALIGN.CENTER)
txt(s, 1, 3.9, 11.3, 0.6, "Descarga la app, inicia sesion y crea tu primer reporte", 18, MUTED, False, PP_ALIGN.CENTER)
txt(s, 1, 5, 11.3, 0.5, "Resumen:", 14, ACCENT, True, PP_ALIGN.CENTER)
txt(s, 1, 5.4, 11.3, 0.5, "Desktop: Panel + Clientes + Equipos + Reportes + Mensajes + Tecnicos", 13, TEXT2, False, PP_ALIGN.CENTER)
txt(s, 1, 5.8, 11.3, 0.5, "Movil: Login + Reportes + Detalle + Crear + Dashboard + Monitoreo + Archivados", 13, TEXT2, False, PP_ALIGN.CENTER)
txt(s, 1, 6.5, 11.3, 0.4, "Dudas? Consulta al gerente", 12, MUTED, False, PP_ALIGN.CENTER)

output = os.path.join(os.path.expanduser("~"), "Desktop", "NexAlert-Capacitacion.pptx")
prs.save(output)
print(f"Capacitacion creada: {output}")
