# -*- coding: utf-8 -*-
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BASE = r"C:\Users\STIVEN\Documents\Default Project"
SSDIR = os.path.join(BASE, "nexalert-screenshots")
LOGO = os.path.join(BASE, "reportes-equipos", "build", "icon-512.png")
OUT = r"C:\Users\STIVEN\Desktop\NexAlert-Capacitacion.pptx"

BG = RGBColor(0x0B, 0x10, 0x1E)
SURF = RGBColor(0x15, 0x1D, 0x35)
SURF2 = RGBColor(0x1E, 0x28, 0x48)
ACC = RGBColor(0x5B, 0x8A, 0xFF)
GRN = RGBColor(0x34, 0xD3, 0x8E)
RED = RGBColor(0xFF, 0x5C, 0x6C)
AMB = RGBColor(0xFF, 0xC1, 0x45)
ORG = RGBColor(0xFF, 0x9A, 0x5C)
PUR = RGBColor(0xB0, 0x7A, 0xFF)
TXT = RGBColor(0xF0, 0xF4, 0xFF)
TXT2 = RGBColor(0xC8, 0xD2, 0xEC)
MUT = RGBColor(0x7E, 0x8A, 0xA8)
WHT = RGBColor(0xFF, 0xFF, 0xFF)
DK = RGBColor(0x18, 0x1C, 0x26)
BLK = RGBColor(0x00, 0x00, 0x00)

FNT = "Segoe UI"
MONO = "Consolas"
LEFT = PP_ALIGN.LEFT
CENTER = PP_ALIGN.CENTER
SW = 13.333

prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def _bg(sl, c):
    sl.background.fill.solid()
    sl.background.fill.fore_color.rgb = c


def _bar(sl):
    for y in (0, prs.slide_height - Inches(0.08)):
        sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, prs.slide_width, Inches(0.08))
        sh.fill.solid()
        sh.fill.fore_color.rgb = ACC
        sh.line.fill.background()
        sh.shadow.inherit = False


def _slide(sl):
    _bg(sl, BG)
    _bar(sl)


def new_slide():
    sl = prs.slides.add_slide(BLANK)
    _slide(sl)
    return sl


def _title(sl, t, sub=None):
    _txt(sl, 0.6, 0.25, 12.1, 0.55, t, sz=26, c=WHT, b=True)
    if sub:
        _txt(sl, 0.6, 0.82, 12.1, 0.35, sub, sz=13, c=MUT)


def _txt(sl, l, t, w, h, text, sz=14, c=TXT, b=False, al=LEFT, fnt=FNT):
    tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = al
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.bold = b
        p.font.name = fnt
    return tb


def _card(sl, l, t, w, h, fill=SURF, line=None):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    else:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _pill(sl, l, t, text, tc, bc, w=None):
    if w is None:
        w = 0.16 * len(text) + 0.4
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(0.34))
    try:
        sh.adjustments[0] = 0.5
    except Exception:
        pass
    sh.fill.solid()
    sh.fill.fore_color.rgb = bc
    sh.line.fill.background()
    sh.shadow.inherit = False
    tf = sh.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.01)
    tf.margin_bottom = Inches(0.01)
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = CENTER
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = tc
    p.font.name = FNT
    return sh


def _circ(sl, l, t, d, fill):
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(d), Inches(d))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _rect(sl, l, t, w, h, fill):
    sh = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _img(sl, path, l, t, w=None, h=None):
    kw = {}
    if w is not None:
        kw["width"] = Inches(w)
    if h is not None:
        kw["height"] = Inches(h)
    return sl.shapes.add_picture(path, Inches(l), Inches(t), **kw)


def _ss(sl, name, l, t, w, h):
    _rect(sl, l + 0.07, t + 0.09, w, h, BLK)
    pic = _img(sl, os.path.join(SSDIR, name), l, t, w=w, h=h)
    fr = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    fr.fill.background()
    fr.line.color.rgb = SURF2
    fr.line.width = Pt(1.5)
    fr.shadow.inherit = False
    return pic


def _phone(sl, l, t):
    body = _card(sl, l, t, 3.2, 5.5, fill=DK, line=SURF2)
    body.line.width = Pt(2)
    _rect(sl, l + 0.12, t + 0.35, 2.96, 4.75, SURF)
    _rect(sl, l + 1.2, t + 0.16, 0.8, 0.05, SURF2)
    _circ(sl, l + 2.15, t + 0.12, 0.12, SURF2)
    _rect(sl, l + 1.1, t + 5.32, 1.0, 0.05, SURF2)
    return body


def _hdr(sl, l, t, w, text, c=ACC):
    _rect(sl, l + 0.25, t + 0.24, 0.07, 0.26, c)
    _txt(sl, l + 0.45, t + 0.18, w - 0.6, 0.35, text, sz=13, c=c, b=True)


def _row(sl, l, t, w, text, c=ACC, sz=11.5, tc=TXT2):
    _rect(sl, l, t + 0.07, 0.07, 0.2, c)
    _txt(sl, l + 0.2, t - 0.04, w - 0.25, 0.32, text, sz=sz, c=tc)


def _feat(sl, l, t, w, h, title, desc, c):
    _card(sl, l, t, w, h, fill=SURF)
    _rect(sl, l, t + 0.12, 0.06, h - 0.24, c)
    _txt(sl, l + 0.18, t + 0.08, w - 0.3, 0.3, title, sz=11.5, c=c, b=True)
    _txt(sl, l + 0.18, t + 0.42, w - 0.3, h - 0.5, desc, sz=9.5, c=TXT2)


def _num_in(shape, n, c):
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(n)
    p.alignment = CENTER
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.name = FNT
    p.font.color.rgb = DK if c in (AMB, GRN, ORG) else WHT


def _stepv(sl, l, t, n, text, c, d=0.42, w=3.0):
    o = _circ(sl, l, t, d, c)
    _num_in(o, n, c)
    _txt(sl, l + d + 0.15, t + 0.03, w, 0.6, text, sz=11, c=TXT2)


def _steph(sl, cx, t, n, label, c):
    o = _circ(sl, cx - 0.25, t, 0.5, c)
    _num_in(o, n, c)
    _txt(sl, cx - 1.05, t + 0.62, 2.1, 0.6, label, sz=9.5, c=TXT2, al=CENTER)


def _arrow(sl, l, t, w=0.4, h=0.22, c=MUT):
    a = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(l), Inches(t), Inches(w), Inches(h))
    a.fill.solid()
    a.fill.fore_color.rgb = c
    a.line.fill.background()
    a.shadow.inherit = False


def s01():
    sl = new_slide()
    _img(sl, LOGO, (SW - 1.5) / 2, 0.75, w=1.5, h=1.5)
    _txt(sl, 0, 2.5, SW, 0.9, "CAPACITACION NEXALERT", sz=48, c=WHT, b=True, al=CENTER)
    _txt(sl, 0, 3.5, SW, 0.4, "Sistema Integral de Gestion de Reportes de Fallas de Equipos", sz=16, c=TXT2, al=CENTER)
    _rect(sl, (SW - 2.5) / 2, 4.3, 2.5, 0.03, ACC)
    _txt(sl, 0, 4.55, SW, 0.35, "Desktop  ·  Movil  ·  Servidor", sz=12, c=MUT, al=CENTER)
    labels = [("Reportes", ACC), ("Sincronizacion", GRN), ("WhatsApp", GRN), ("Offline", AMB), ("Push", PUR), ("Biometria", ORG)]
    ws = [0.16 * len(s) + 0.4 for s, _ in labels]
    x = (SW - (sum(ws) + 0.25 * (len(ws) - 1))) / 2
    for (s, col), w in zip(labels, ws):
        _pill(sl, x, 6.1, s, WHT, col, w=w)
        x += w + 0.25


def s02():
    sl = new_slide()
    _title(sl, "QUE ES NEXALERT", "Plataforma completa de reportes de fallas multi-plataforma")
    _ss(sl, "01-panel.png", 0.6, 1.4, 7.8, 5.2)
    feats = [
        ("Desktop Electron", "App nativa para Windows", ACC),
        ("App Movil", "Android APK offline-first", GRN),
        ("Servidor", "API REST + PostgreSQL", PUR),
        ("WhatsApp", "Notificaciones automaticas", GRN),
        ("Firebase Push", "Alertas en tiempo real", ORG),
        ("Sincronizacion", "Cola local con reintentos", AMB),
        ("Auto-update", "Actualizador integrado", RED),
        ("Biometria", "Huella digital en el movil", ACC),
    ]
    for i, (t_, d_, c_) in enumerate(feats):
        col = i % 2
        rw = i // 2
        _feat(sl, 8.7 + col * 2.18, 1.4 + rw * 1.34, 2.02, 1.22, t_, d_, c_)


def s03():
    sl = new_slide()
    _title(sl, "ARQUITECTURA DEL SISTEMA", "Tres componentes que trabajan en conjunto")
    arch = [
        (0.6, "DESKTOP ELECTRON", ACC, ["Interfaz grafica HTML / CSS / JS", "Backend local con Node.js", "Base de datos SQLite local", "Generador de reportes PDF", "Auto-actualizador integrado", "Accesos por rol de usuario"]),
        (4.77, "APP MOVIL ANDROID", GRN, ["Instalacion por APK directo", "Base local SQLite", "Camara para evidencias", "GPS en cada reporte", "Modo offline completo", "Login con huella digital"]),
        (8.94, "SERVIDOR CLOUD", PUR, ["API REST Node.js / Express", "PostgreSQL centralizado", "Autenticacion con tokens JWT", "Almacenamiento de imagenes", "Integracion WhatsApp", "Backup diario automatico"]),
    ]
    for x, tt, cc, items in arch:
        _card(sl, x, 1.35, 3.79, 3.25, fill=SURF)
        _hdr(sl, x, 1.35, 3.79, tt, cc)
        for j, it in enumerate(items):
            _row(sl, x + 0.3, 2.0 + j * 0.42, 3.4, it, cc, sz=10.5)
    _card(sl, 0.6, 4.85, 12.13, 2.35, fill=SURF)
    _hdr(sl, 0.6, 4.85, 5.0, "FLUJO DE DATOS", ACC)
    steps = [("Captura del reporte", ACC), ("Cola local offline", AMB), ("Sync via API REST", GRN), ("PostgreSQL central", PUR), ("WhatsApp + Push", ORG)]
    xs = [1.7, 4.15, 6.6, 9.05, 11.5]
    for k, (lab, cc) in enumerate(steps):
        _steph(sl, xs[k], 5.5, k + 1, lab, cc)
        if k < 4:
            _arrow(sl, xs[k] + 0.45, 5.63, w=1.3, h=0.24)


def s04():
    sl = new_slide()
    _title(sl, "ROLES Y PERMISOS", "Cada usuario ve solo lo que necesita")
    _card(sl, 0.6, 1.35, 5.95, 5.75, fill=SURF)
    _hdr(sl, 0.6, 1.35, 5.95, "GERENTE", ACC)
    _txt(sl, 0.9, 1.85, 5.3, 0.3, "Acceso total al sistema", sz=10.5, c=MUT)
    gerente = [("Panel general", "KPIs y estadisticas globales"), ("Reportes", "Control completo del ciclo"), ("Clientes", "Altas, bajas y edicion"), ("Equipos", "Inventario maestro"), ("Tecnicos", "Gestion del personal"), ("Mensajes WhatsApp", "Historial y envios")]
    for j, (tt, dd) in enumerate(gerente):
        y = 2.45 + j * 0.62
        _rect(sl, 0.9, y + 0.05, 0.07, 0.3, ACC)
        _txt(sl, 1.1, y - 0.05, 2.3, 0.3, tt, sz=12.5, c=WHT, b=True)
        _txt(sl, 3.3, y, 3.1, 0.3, dd, sz=10, c=MUT)
    _card(sl, 6.78, 1.35, 5.95, 5.75, fill=SURF)
    _hdr(sl, 6.78, 1.35, 5.95, "TECNICO", GRN)
    _txt(sl, 7.08, 1.85, 5.3, 0.3, "Acceso operativo en campo", sz=10.5, c=MUT)
    tecnico = ["Mis reportes asignados", "Nuevo reporte", "Detalle y cambio de estado", "Fotos de evidencia", "Ubicacion GPS", "Consulta de clientes", "Consulta de equipos", "Mensajes WhatsApp", "Mi perfil", "Notificaciones push", "Modo offline"]
    for j, tt in enumerate(tecnico):
        y = 2.35 + j * 0.42
        _rect(sl, 7.08, y + 0.07, 0.07, 0.2, GRN)
        _txt(sl, 7.28, y - 0.04, 5.2, 0.3, tt, sz=11.5, c=TXT2)


def s05():
    sl = new_slide()
    _title(sl, "PANEL PRINCIPAL DESKTOP", "Vista general del sistema en tiempo real")
    _ss(sl, "01-panel.png", 0.6, 1.4, 8.0, 5.2)
    _card(sl, 8.85, 1.4, 3.88, 5.2, fill=SURF)
    _hdr(sl, 8.85, 1.4, 3.88, "INDICADORES CLAVE", ACC)
    stats = [("128", "Reportes totales", ACC), ("23", "Pendientes hoy", AMB), ("17", "En proceso", ORG), ("342", "Resueltos del mes", GRN), ("8", "Tecnicos activos", PUR), ("56", "Clientes registrados", ACC), ("214", "Equipos inventariados", GRN)]
    for j, (v, lab, cc) in enumerate(stats):
        y = 2.0 + j * 0.52
        _txt(sl, 9.1, y, 0.85, 0.35, v, sz=17, c=cc, b=True)
        _txt(sl, 9.95, y + 0.08, 2.6, 0.3, lab, sz=10, c=MUT)
    pills = [("Por tecnico", ACC), ("Por cliente", PUR), ("Por estado", ORG), ("Por fecha", GRN)]
    for j, (lab, cc) in enumerate(pills):
        px = 9.1 + (j % 2) * 1.85
        py = 5.95 + (j // 2) * 0.42
        _pill(sl, px, py, lab, DK if cc in (AMB, GRN, ORG) else WHT, cc, w=1.7)


def s06():
    sl = new_slide()
    _title(sl, "MODULO REPORTES", "Busqueda avanzada y gestion completa")
    _ss(sl, "02-reportes.png", 0.6, 1.4, 7.8, 5.2)
    _card(sl, 8.65, 1.4, 4.08, 5.2, fill=SURF)
    _hdr(sl, 8.65, 1.4, 4.08, "FILTROS DISPONIBLES", ACC)
    filtros = ["Busqueda libre por texto", "Estado del reporte", "Nivel de prioridad", "Cliente", "Equipo", "Tecnico asignado", "Rango de fechas", "Tipo de problema", "Orden de resultados"]
    for j, f in enumerate(filtros):
        _row(sl, 8.95, 1.95 + j * 0.33, 3.6, f, ACC, sz=10.5)
    _hdr(sl, 8.65, 4.95, 4.08, "ACCIONES", GRN)
    acciones = [("Ver detalle", ACC), ("Cambiar estado", AMB), ("Exportar Excel", GRN), ("Imprimir PDF", PUR), ("Asignar tecnico", ORG), ("Eliminar", RED)]
    for j, (lab, cc) in enumerate(acciones):
        px = 8.95 + (j % 2) * 1.95
        py = 5.5 + (j // 2) * 0.42
        _pill(sl, px, py, lab, DK if cc in (AMB, GRN, ORG) else WHT, cc, w=1.8)


def s07():
    sl = new_slide()
    _title(sl, "CREAR REPORTE DESKTOP", "Formulario guiado paso a paso")
    _ss(sl, "03-nuevo-reporte.png", 0.6, 1.4, 7.8, 5.2)
    _card(sl, 8.65, 1.4, 4.08, 5.2, fill=SURF)
    _hdr(sl, 8.65, 1.4, 4.08, "CAMPOS DEL FORMULARIO", ACC)
    campos = [("Cliente", "lista desplegable"), ("Equipo", "filtrado por cliente"), ("Descripcion del problema", "texto libre"), ("Prioridad", "Alta / Media / Baja"), ("Tecnico asignado", "opcional"), ("Fecha y hora", "automatica")]
    for j, (tt, dd) in enumerate(campos):
        y = 1.95 + j * 0.4
        _rect(sl, 8.95, y + 0.07, 0.07, 0.2, ACC)
        _txt(sl, 9.15, y - 0.04, 1.95, 0.3, tt, sz=11, c=WHT, b=True)
        _txt(sl, 11.1, y, 1.5, 0.3, dd, sz=9, c=MUT)
    _hdr(sl, 8.65, 4.5, 4.08, "ESTADOS POSIBLES", ORG)
    estados = [("PENDIENTE", AMB), ("EN PROCESO", ACC), ("ESPERA REPUESTO", ORG), ("RESUELTO", GRN), ("CERRADO", MUT)]
    for j, (lab, cc) in enumerate(estados):
        _pill(sl, 8.95, 4.95 + j * 0.42, lab, DK if cc in (AMB, GRN, ORG) else WHT, cc, w=2.2)


def s08():
    sl = new_slide()
    _title(sl, "CLIENTES Y EQUIPOS", "Maestros principales del sistema")
    _pill(sl, 0.6, 1.28, "CLIENTES", WHT, ACC, w=1.5)
    _pill(sl, 6.85, 1.28, "EQUIPOS", WHT, GRN, w=1.5)
    _ss(sl, "04-clientes.png", 0.6, 1.75, 5.95, 4.9)
    _ss(sl, "06-equipos.png", 6.85, 1.75, 5.95, 4.9)


def s09():
    sl = new_slide()
    _title(sl, "GESTION DE TECNICOS", "Personal operativo del sistema")
    _ss(sl, "05-tecnicos.png", 0.6, 1.4, 7.8, 5.2)
    _card(sl, 8.65, 1.4, 4.08, 5.2, fill=SURF)
    _hdr(sl, 8.65, 1.4, 4.08, "DATOS DEL TECNICO", ACC)
    datos = ["Nombre completo", "Usuario de acceso", "Contrasena inicial", "Telefono WhatsApp", "Especialidad tecnica", "Estado activo / inactivo"]
    for j, d in enumerate(datos):
        _row(sl, 8.95, 1.95 + j * 0.36, 3.6, d, ACC, sz=11)
    _hdr(sl, 8.65, 4.3, 4.08, "FUNCIONES", GRN)
    funcs = ["Crear nuevo tecnico", "Editar informacion", "Activar / desactivar cuenta", "Restablecer contrasena", "Ver carga de trabajo"]
    for j, f in enumerate(funcs):
        _row(sl, 8.95, 4.75 + j * 0.44, 3.6, f, GRN, sz=11)


def s10():
    sl = new_slide()
    _title(sl, "MENSAJES WHATSAPP", "Comunicacion directa con los clientes")
    _ss(sl, "07-mensajes.png", 0.6, 1.4, 7.8, 5.2)
    _card(sl, 8.65, 1.4, 4.08, 5.2, fill=SURF)
    _hdr(sl, 8.65, 1.4, 4.08, "FUNCIONES DE CHAT", GRN)
    funcs = ["Historial por cliente", "Envio manual de mensajes", "Plantillas rapidas", "Adjuntar fotografias", "Estado enviado / leido", "Respuestas del cliente", "Busqueda en conversaciones", "Copiar mensaje", "Reenviar reporte completo", "Numeros verificados", "Registro de envios (log)", "Enlace directo al reporte"]
    for j, f in enumerate(funcs):
        _row(sl, 8.95, 1.95 + j * 0.4, 3.6, f, GRN, sz=10.5)


def s11():
    sl = new_slide()
    _title(sl, "LOGIN Y BIOMETRIA", "Acceso seguro en la aplicacion movil")
    l, t = 0.7, 1.35
    _phone(sl, l, t)
    sx, sy = l + 0.12, t + 0.35
    _txt(sl, sx, sy + 0.2, 2.96, 0.4, "NexAlert", sz=19, c=ACC, b=True, al=CENTER)
    for yy, lab in ((sy + 0.8, "Usuario"), (sy + 1.37, "Contrasena")):
        f = _rect(sl, sx + 0.25, yy, 2.46, 0.42, SURF2)
        pf = f.text_frame.paragraphs[0]
        pf.text = lab
        pf.alignment = CENTER
        pf.font.size = Pt(10.5)
        pf.font.color.rgb = MUT
        pf.font.name = FNT
    btn = _rect(sl, sx + 0.25, sy + 2.0, 2.46, 0.46, ACC)
    pb = btn.text_frame.paragraphs[0]
    pb.text = "Entrar"
    pb.alignment = CENTER
    pb.font.size = Pt(12)
    pb.font.bold = True
    pb.font.color.rgb = WHT
    pb.font.name = FNT
    bio = _rect(sl, sx + 0.25, sy + 2.65, 2.46, 0.46, SURF2)
    bio.line.color.rgb = PUR
    bio.line.width = Pt(1)
    pbio = bio.text_frame.paragraphs[0]
    pbio.text = "Usar huella digital"
    pbio.alignment = CENTER
    pbio.font.size = Pt(11)
    pbio.font.bold = True
    pbio.font.color.rgb = PUR
    pbio.font.name = FNT
    _circ(sl, sx + 1.31, sy + 3.4, 0.48, PUR)
    _card(sl, 4.15, 1.4, 4.3, 5.6, fill=SURF)
    _hdr(sl, 4.15, 1.4, 4.3, "FLUJO DE AUTENTICACION", ACC)
    pasos = ["Abrir la aplicacion movil", "Huella digital o PIN de respaldo", "Validacion local cifrada", "Acceso directo sin escribir clave"]
    for j, p_ in enumerate(pasos):
        _stepv(sl, 4.45, 2.2 + j * 1.15, j + 1, p_, ACC, w=3.2)
    _card(sl, 8.7, 1.4, 4.03, 5.6, fill=SURF)
    _hdr(sl, 8.7, 1.4, 4.03, "SEGURIDAD IMPLEMENTADA", PUR)
    sec = ["Contrasenas con hash bcrypt", "Token JWT expira en 24 horas", "Huella guardada solo en el dispositivo", "Base local cifrada", "Bloqueo tras 5 intentos fallidos", "Roles con permisos separados", "Cierre de sesion remoto"]
    for j, s_ in enumerate(sec):
        _row(sl, 9.0, 2.15 + j * 0.62, 3.5, s_, PUR, sz=11)


def s12():
    sl = new_slide()
    _title(sl, "DETALLE DE REPORTE EN EL MOVIL", "Toda la informacion en campo")
    l, t = 0.7, 1.35
    _phone(sl, l, t)
    sx, sy = l + 0.12, t + 0.35
    hd = _rect(sl, sx, sy, 2.96, 0.5, SURF2)
    ph = hd.text_frame.paragraphs[0]
    ph.text = "REPORTE #0042"
    ph.alignment = CENTER
    ph.font.size = Pt(11)
    ph.font.bold = True
    ph.font.color.rgb = WHT
    ph.font.name = FNT
    _txt(sl, sx, sy + 0.58, 2.96, 0.25, "12/08/2026  ·  09:41", sz=9.5, c=MUT, al=CENTER)
    _txt(sl, sx, sy + 0.85, 2.96, 0.35, "Impresora HP no enciende", sz=12, c=WHT, b=True, al=CENTER)
    _pill(sl, sx + 0.8, sy + 1.3, "EN PROCESO", DK, AMB, w=1.35)
    _txt(sl, sx, sy + 1.78, 2.96, 0.25, "GPS: -12.046, -77.043", sz=9.5, c=PUR, al=CENTER)
    btns = [("Llamar", GRN), ("WhatsApp", ACC), ("Foto", PUR), ("Resolver", ORG)]
    for j, (lab, cc) in enumerate(btns):
        bx = sx + 0.25 + (j % 2) * 1.36
        by = sy + 2.2 + (j // 2) * 0.52
        b = _rect(sl, bx, by, 1.3, 0.42, cc)
        tb = b.text_frame.paragraphs[0]
        tb.text = lab
        tb.alignment = CENTER
        tb.font.size = Pt(10)
        tb.font.bold = True
        tb.font.color.rgb = DK if cc in (AMB, GRN, ORG) else WHT
        tb.font.name = FNT
    cm = _rect(sl, sx + 0.25, sy + 3.35, 2.46, 0.6, SURF2)
    tc = cm.text_frame.paragraphs[0]
    tc.text = "Agregar comentario..."
    tc.font.size = Pt(9)
    tc.font.color.rgb = MUT
    tc.font.name = FNT
    _card(sl, 4.15, 1.4, 4.3, 5.6, fill=SURF)
    _hdr(sl, 4.15, 1.4, 4.3, "CAMPOS DEL REPORTE (13)", ACC)
    campos = ["ID unico del reporte", "Fecha y hora de creacion", "Cliente asociado", "Equipo afectado", "Modelo y serie", "Descripcion del problema", "Nivel de prioridad", "Estado actual", "Tecnico asignado", "Fotografia de evidencia", "Coordenadas GPS", "Ultima actualizacion", "Comentarios tecnicos"]
    for j, c_ in enumerate(campos):
        _row(sl, 4.45, 2.0 + j * 0.37, 3.8, c_, ACC, sz=10.5)
    _card(sl, 8.7, 1.4, 4.03, 5.6, fill=SURF)
    _hdr(sl, 8.7, 1.4, 4.03, "ACCIONES DISPONIBLES (9)", GRN)
    acc = ["Cambiar estado del reporte", "Agregar comentario tecnico", "Tomar foto de evidencia", "Ver foto en pantalla completa", "Abrir ubicacion en el mapa", "Llamar al cliente", "Enviar WhatsApp al cliente", "Marcar como resuelto", "Volver a la lista"]
    for j, a_ in enumerate(acc):
        _row(sl, 9.0, 2.05 + j * 0.55, 3.5, a_, GRN, sz=11)


def s13():
    sl = new_slide()
    _title(sl, "CREAR REPORTE EN EL MOVIL", "Alta de reportes desde campo en segundos")
    l, t = 0.7, 1.35
    _phone(sl, l, t)
    sx, sy = l + 0.12, t + 0.35
    _txt(sl, sx, sy + 0.12, 2.96, 0.35, "NUEVO REPORTE", sz=14, c=ACC, b=True, al=CENTER)
    for yy, lab in ((sy + 0.6, "Cliente  *"), (sy + 1.15, "Equipo  *")):
        f = _rect(sl, sx + 0.25, yy, 2.46, 0.4, SURF2)
        pf = f.text_frame.paragraphs[0]
        pf.text = lab
        pf.font.size = Pt(10)
        pf.font.color.rgb = TXT2
        pf.font.name = FNT
    prb = _rect(sl, sx + 0.25, sy + 1.7, 2.46, 0.75, SURF2)
    pp = prb.text_frame.paragraphs[0]
    pp.text = "Descripcion del problema  *"
    pp.font.size = Pt(9.5)
    pp.font.color.rgb = MUT
    pp.font.name = FNT
    chips = [("ALTA", RED, 0.72), ("MEDIA", AMB, 0.86), ("BAJA", GRN, 0.72)]
    cx = sx + 0.25
    for lab, cc, cw in chips:
        _pill(sl, cx, sy + 2.6, lab, DK if cc in (AMB, GRN) else WHT, cc, w=cw)
        cx += cw + 0.06
    fb = _rect(sl, sx + 0.25, sy + 3.15, 1.18, 0.4, PUR)
    tf1 = fb.text_frame.paragraphs[0]
    tf1.text = "Tomar foto"
    tf1.alignment = CENTER
    tf1.font.size = Pt(9.5)
    tf1.font.bold = True
    tf1.font.color.rgb = WHT
    tf1.font.name = FNT
    gb = _rect(sl, sx + 1.53, sy + 3.15, 1.18, 0.4, ACC)
    tf2 = gb.text_frame.paragraphs[0]
    tf2.text = "GPS"
    tf2.alignment = CENTER
    tf2.font.size = Pt(9.5)
    tf2.font.bold = True
    tf2.font.color.rgb = WHT
    tf2.font.name = FNT
    sv = _rect(sl, sx + 0.25, sy + 3.75, 2.46, 0.45, GRN)
    ts = sv.text_frame.paragraphs[0]
    ts.text = "GUARDAR"
    ts.alignment = CENTER
    ts.font.size = Pt(12)
    ts.font.bold = True
    ts.font.color.rgb = DK
    ts.font.name = FNT
    obl = [(4.15, "CLIENTE  *", "Seleccion obligatoria de la lista"), (7.05, "EQUIPO  *", "Depende del cliente elegido"), (9.95, "PROBLEMA  *", "Minimo 10 caracteres")]
    for x, tt, dd in obl:
        _card(sl, x, 1.4, 2.75, 1.45, fill=SURF, line=RED)
        _txt(sl, x + 0.2, 1.55, 2.35, 0.3, tt, sz=12, c=RED, b=True)
        _txt(sl, x + 0.2, 1.95, 2.35, 0.7, dd, sz=9.5, c=TXT2)
    _card(sl, 4.15, 3.05, 8.55, 3.95, fill=SURF)
    _hdr(sl, 4.15, 3.05, 8.55, "FUNCIONES DEL FORMULARIO", GRN)
    funcs = ["Validacion de campos vacios", "Lista de clientes disponible offline", "Foto comprimida automaticamente", "GPS capturado al abrir la app", "Borrador se guarda solo", "Confirmacion antes de guardar"]
    cols = [ACC, GRN, PUR, ORG, AMB, RED]
    for j, f_ in enumerate(funcs):
        fx = 4.5 + (j % 2) * 4.2
        fy = 3.65 + (j // 2) * 0.85
        _rect(sl, fx, fy + 0.07, 0.07, 0.24, cols[j])
        _txt(sl, fx + 0.2, fy, 3.7, 0.55, f_, sz=11.5, c=TXT2)


def s14():
    sl = new_slide()
    _title(sl, "ESTADOS Y FLUJO DE TRABAJO", "Ciclo de vida completo del reporte")
    _card(sl, 0.6, 1.35, 12.13, 1.9, fill=SURF)
    states = [("PENDIENTE", AMB), ("EN PROCESO", ACC), ("ESPERA REPUESTO", ORG), ("RESUELTO", GRN), ("CERRADO", MUT)]
    xs = [1.0, 3.35, 5.7, 8.05, 10.4]
    for k, (lab, cc) in enumerate(states):
        b = _card(sl, xs[k], 2.05, 1.95, 0.75, fill=SURF2)
        _rect(sl, xs[k], 2.05, 0.07, 0.75, cc)
        pb = b.text_frame.paragraphs[0]
        pb.text = lab
        pb.alignment = CENTER
        pb.font.size = Pt(11)
        pb.font.bold = True
        pb.font.color.rgb = cc
        pb.font.name = FNT
        if k < 4:
            _arrow(sl, xs[k] + 2.0, 2.31, w=0.3, h=0.24)
    _card(sl, 0.6, 3.5, 6.4, 3.6, fill=SURF)
    _hdr(sl, 0.6, 3.5, 6.4, "TRANSICIONES PERMITIDAS", ACC)
    trans = ["Pendiente -> En proceso", "En proceso -> Espera repuesto", "Espera repuesto -> En proceso", "En proceso -> Resuelto", "Resuelto -> Cerrado", "Resuelto -> Reabrir (pendiente)", "Cancelar (solo gerente)"]
    for j, tr in enumerate(trans):
        _pill(sl, 0.9, 4.15 + j * 0.42, tr, TXT2, SURF2, w=5.6)
    _card(sl, 7.2, 3.5, 5.53, 3.6, fill=SURF)
    _hdr(sl, 7.2, 3.5, 5.53, "PRIORIDADES", RED)
    prios = [("ALTA", "Atencion inmediata", RED), ("MEDIA", "Se atiende el mismo dia", AMB), ("BAJA", "Resolucion esta semana", GRN), ("PROGRAMADA", "Mantenimiento planificado", PUR)]
    for j, (tt, dd, cc) in enumerate(prios):
        y = 4.15 + j * 0.72
        _card(sl, 7.5, y, 4.9, 0.6, fill=SURF2)
        _rect(sl, 7.5, y, 0.07, 0.6, cc)
        _txt(sl, 7.7, y + 0.13, 1.7, 0.3, tt, sz=11.5, c=cc, b=True)
        _txt(sl, 9.45, y + 0.15, 2.85, 0.3, dd, sz=10, c=TXT2)


def s15():
    sl = new_slide()
    _title(sl, "INTEGRACION WHATSAPP", "Notificaciones automaticas al cliente")
    _card(sl, 0.6, 1.35, 3.9, 5.75, fill=SURF)
    _hdr(sl, 0.6, 1.35, 3.9, "CONEXION (5 PASOS)", GRN)
    pasos = ["Escanear codigo QR", "La sesion queda vinculada", "Credenciales guardadas", "Reconexion automatica", "Alerta si se desconecta"]
    for j, p_ in enumerate(pasos):
        _stepv(sl, 0.9, 2.15 + j * 0.95, j + 1, p_, GRN, w=2.9)
    _card(sl, 4.7, 1.35, 4.0, 5.75, fill=SURF)
    _hdr(sl, 4.7, 1.35, 4.0, "FORMATO DEL MENSAJE", ACC)
    _rect(sl, 4.95, 2.0, 3.5, 3.4, DK)
    msg = "*NEXALERT - REPORTE #0042*\n\nCliente: Clinica Central\nEquipo: Impresora HP 4020\nProblema: No enciende\nPrioridad: ALTA\nEstado: EN PROCESO\nTecnico: J. Perez\n\nFecha: 12/08/2026 09:41"
    _txt(sl, 5.15, 2.15, 3.1, 3.1, msg, sz=10, c=GRN, fnt=MONO)
    _txt(sl, 4.95, 5.6, 3.5, 0.6, "Formato editable desde el panel de configuracion", sz=9.5, c=MUT)
    _card(sl, 8.9, 1.35, 3.83, 5.75, fill=SURF)
    _hdr(sl, 8.9, 1.35, 3.83, "FUNCIONES (10)", PUR)
    funcs = ["Envio automatico al crear reporte", "Aviso en cada cambio de estado", "Resumen al resolver", "Plantillas configurables", "Historial por numero", "Adjuntar foto de evidencia", "Confirmacion de lectura", "Reenvio manual", "Multiples numeros por cliente", "Log de errores de envio"]
    for j, f_ in enumerate(funcs):
        _row(sl, 9.2, 2.0 + j * 0.47, 3.3, f_, PUR, sz=10.5)


def s16():
    sl = new_slide()
    _title(sl, "MODO OFFLINE Y SINCRONIZACION", "El sistema nunca se detiene")
    _card(sl, 0.6, 1.35, 5.95, 2.9, fill=SURF)
    _hdr(sl, 0.6, 1.35, 5.95, "FLUJO OFFLINE (5 PASOS)", AMB)
    pasos = ["La app sigue funcionando sin red", "Los datos se guardan en SQLite local", "Cada cambio entra a la cola de sync", "Al detectar red reintenta solo", "Orden FIFO garantiza consistencia"]
    for j, p_ in enumerate(pasos):
        _stepv(sl, 0.9, 1.95 + j * 0.42, j + 1, p_, AMB, d=0.32, w=4.9)
    _card(sl, 6.78, 1.35, 5.95, 2.9, fill=SURF)
    _hdr(sl, 6.78, 1.35, 5.95, "COLA DE SINCRONIZACION", GRN)
    eps = ["POST /api/reportes   -   crear", "PUT  /api/reportes/:id   -   actualizar", "POST /api/clientes   -   nuevo cliente", "POST /api/equipos   -   nuevo equipo", "POST /api/mensajes   -   chat"]
    for j, e in enumerate(eps):
        _txt(sl, 7.08, 1.95 + j * 0.42, 5.4, 0.32, e, sz=10.5, c=GRN, fnt=MONO)
    _card(sl, 0.6, 4.45, 12.13, 2.65, fill=SURF)
    _hdr(sl, 0.6, 4.45, 12.13, "DETALLES TECNICOS", PUR)
    det1 = ["Timestamps UTC en cada registro", "IDs UUID evitan duplicados", "Reintentos con backoff exponencial", "Indicador visual online / offline"]
    det2 = ["Compresion de imagenes antes de subir", "Conflictos: ultima escritura gana", "Login offline validado con hash local"]
    for j, d_ in enumerate(det1):
        _row(sl, 0.95, 5.05 + j * 0.5, 5.6, d_, PUR, sz=11.5)
    for j, d_ in enumerate(det2):
        _row(sl, 6.9, 5.05 + j * 0.5, 5.6, d_, PUR, sz=11.5)


def s17():
    sl = new_slide()
    _title(sl, "MONITOREO DE TECNICOS", "Supervision en tiempo real")
    l, t = 0.7, 1.35
    _phone(sl, l, t)
    sx, sy = l + 0.12, t + 0.35
    _txt(sl, sx, sy + 0.12, 2.96, 0.35, "MONITOREO", sz=14, c=ACC, b=True, al=CENTER)
    minis = [("5", "ACTIVOS", GRN), ("23", "PENDIENTES", AMB)]
    for j, (v, lab, cc) in enumerate(minis):
        mx = sx + 0.25 + j * 1.28
        _card(sl, mx, sy + 0.6, 1.18, 0.7, fill=SURF2)
        _txt(sl, mx, sy + 0.63, 1.18, 0.35, v, sz=16, c=cc, b=True, al=CENTER)
        _txt(sl, mx, sy + 1.0, 1.18, 0.25, lab, sz=8, c=MUT, al=CENTER)
    tec = [("J. Perez", "4 trabajos", GRN), ("M. Lopez", "3 trabajos", GRN), ("C. Ruiz", "1 trabajo", AMB), ("A. Gomez", "offline", RED), ("L. Torres", "2 trabajos", GRN)]
    for j, (nm, st, cc) in enumerate(tec):
        y = sy + 1.55 + j * 0.62
        _circ(sl, sx + 0.25, y + 0.06, 0.14, cc)
        _txt(sl, sx + 0.48, y - 0.04, 1.4, 0.3, nm, sz=10, c=WHT, b=True)
        _txt(sl, sx + 1.85, y, 1.0, 0.3, st, sz=9, c=MUT)
    _card(sl, 4.15, 1.4, 4.3, 5.6, fill=SURF)
    _hdr(sl, 4.15, 1.4, 4.3, "ESTADISTICAS (7)", ACC)
    est = ["Reportes creados por dia", "Tiempo promedio de resolucion", "Tecnicos conectados ahora", "Distribucion por estado", "Carga de trabajo por tecnico", "Clientes con mas reportes", "Equipos con fallas repetidas"]
    for j, e in enumerate(est):
        _row(sl, 4.45, 2.1 + j * 0.62, 3.8, e, ACC, sz=11)
    _card(sl, 8.7, 1.4, 4.03, 5.6, fill=SURF)
    _hdr(sl, 8.7, 1.4, 4.03, "ACCIONES DEL GERENTE (7)", ORG)
    acc = ["Reasignar tecnico de un reporte", "Escalar prioridad a ALTA", "Enviar recordatorio por WhatsApp", "Ver mapa de ubicaciones", "Exportar estadisticas a Excel", "Crear backup manual", "Configurar alertas automaticas"]
    for j, a_ in enumerate(acc):
        _row(sl, 9.0, 2.1 + j * 0.62, 3.5, a_, ORG, sz=11)


def s18():
    sl = new_slide()
    _title(sl, "NOTIFICACIONES PUSH", "Firebase Cloud Messaging")
    _card(sl, 0.6, 1.35, 3.9, 5.75, fill=SURF)
    _hdr(sl, 0.6, 1.35, 3.9, "EVENTOS (6)", ACC)
    evs = [("Nuevo reporte asignado", ACC), ("Cambio de estado", AMB), ("Reporte resuelto", GRN), ("Recordatorio de pendientes", ORG), ("Servidor sin conexion", RED), ("Backup completado", PUR)]
    for j, (e, cc) in enumerate(evs):
        _row(sl, 0.9, 2.1 + j * 0.78, 3.3, e, cc, sz=11.5)
    _card(sl, 4.7, 1.35, 4.0, 5.75, fill=SURF)
    _hdr(sl, 4.7, 1.35, 4.0, "CONFIGURACION TECNICA (7)", GRN)
    cfg = ["Firebase Cloud Messaging", "Token unico por dispositivo", "Topics por rol (gerente / tecnico)", "Prioridad alta en Android", "Icono y color corporativos", "Deep link abre el reporte", "Reintento si falla la entrega"]
    for j, c_ in enumerate(cfg):
        _row(sl, 5.0, 2.05 + j * 0.68, 3.4, c_, GRN, sz=11)
    _card(sl, 8.9, 1.35, 3.83, 5.75, fill=SURF)
    _hdr(sl, 8.9, 1.35, 3.83, "EJEMPLOS", PUR)
    notis = [
        ("NexAlert  ·  ahora", "Nuevo reporte #0042", "Impresora HP no enciende - Clinica Central (ALTA)"),
        ("NexAlert  ·  hace 5 min", "Reporte #0038 RESUELTO", "J. Perez finalizo el mantenimiento del equipo"),
        ("NexAlert  ·  hace 1 hora", "Prioridad escalada", "UPS de Sala de Servidores sin bateria - ver ahora"),
    ]
    for j, (tt, bt, bd) in enumerate(notis):
        y = 2.0 + j * 1.55
        _card(sl, 9.1, y, 3.43, 1.35, fill=DK, line=SURF2)
        _circ(sl, 9.3, y + 0.18, 0.12, ACC)
        _txt(sl, 9.52, y + 0.1, 2.9, 0.25, tt, sz=9, c=MUT)
        _txt(sl, 9.3, y + 0.42, 3.05, 0.28, bt, sz=10.5, c=WHT, b=True)
        _txt(sl, 9.3, y + 0.74, 3.05, 0.55, bd, sz=9, c=TXT2)


def s19():
    sl = new_slide()
    _title(sl, "PROBLEMAS COMUNES Y SOLUCIONES")
    _txt(sl, 0.75, 1.28, 3.4, 0.3, "PROBLEMA", sz=11, c=MUT, b=True)
    _txt(sl, 4.35, 1.28, 4.0, 0.3, "CAUSA PROBABLE", sz=11, c=MUT, b=True)
    _txt(sl, 8.55, 1.28, 4.0, 0.3, "SOLUCION", sz=11, c=GRN, b=True)
    rows = [
        ("No genera el PDF", "Impresora no predeterminada", "Elegir impresora en Configuracion"),
        ("La app no sincroniza", "Sin internet o servidor caido", "Verificar indicador online y reintentar"),
        ("WhatsApp no envia", "Sesion QR expirada", "Reescanear QR en el panel Mensajes"),
        ("La huella no funciona", "Huella cambiada en el telefono", "Ingresar con clave y re-registrar"),
        ("No se adjunta la foto", "Permiso de camara denegado", "Habilitar permisos en Ajustes Android"),
        ("Login rechazado", "Clave vencida o mal escrita", "Solicitar restablecimiento al gerente"),
        ("El panel va lento", "Demasiados registros sin filtro", "Aplicar filtros y rango de fechas"),
        ("No llega el push", "Token invalido tras reinstalar", "Abrir la app para regenerar el token"),
    ]
    for j, (pb, cz, sol) in enumerate(rows):
        y = 1.7 + j * 0.66
        _card(sl, 0.6, y, 12.13, 0.6, fill=SURF if j % 2 == 0 else SURF2)
        _txt(sl, 0.85, y + 0.14, 3.4, 0.32, pb, sz=11, c=WHT, b=True)
        _txt(sl, 4.35, y + 0.16, 4.0, 0.32, cz, sz=10.5, c=TXT2)
        _txt(sl, 8.55, y + 0.16, 4.0, 0.32, sol, sz=10.5, c=GRN)


def s20():
    sl = new_slide()
    _img(sl, LOGO, (SW - 1.4) / 2, 0.6, w=1.4, h=1.4)
    _txt(sl, 0, 2.35, SW, 0.8, "GRACIAS POR SU ATENCION", sz=44, c=WHT, b=True, al=CENTER)
    _rect(sl, (SW - 2.5) / 2, 3.35, 2.5, 0.03, ACC)
    _txt(sl, 0, 3.55, SW, 0.35, "NexAlert - Reportes de fallas de equipos", sz=14, c=MUT, al=CENTER)
    _card(sl, 1.6, 4.35, 4.9, 2.6, fill=SURF)
    _hdr(sl, 1.6, 4.35, 4.9, "MODULOS", ACC)
    mods = ["Panel Desktop", "Reportes", "Clientes", "Equipos", "Tecnicos", "Mensajes WhatsApp", "App Movil", "Monitoreo"]
    for j, m in enumerate(mods):
        mx = 1.95 + (j % 2) * 2.25
        my = 5.05 + (j // 2) * 0.42
        _rect(sl, mx, my + 0.07, 0.07, 0.2, ACC)
        _txt(sl, mx + 0.18, my - 0.04, 2.05, 0.3, m, sz=10.5, c=TXT2)
    _card(sl, 6.93, 4.35, 4.9, 2.6, fill=SURF)
    _hdr(sl, 6.93, 4.35, 4.9, "TECNOLOGIAS", GRN)
    techs = ["Electron", "Node.js", "SQLite", "PostgreSQL", "Express", "Firebase FCM", "WhatsApp Web", "Android APK"]
    for j, m in enumerate(techs):
        mx = 7.28 + (j % 2) * 2.25
        my = 5.05 + (j // 2) * 0.42
        _rect(sl, mx, my + 0.07, 0.07, 0.2, GRN)
        _txt(sl, mx + 0.18, my - 0.04, 2.05, 0.3, m, sz=10.5, c=TXT2)


if __name__ == "__main__":
    for fn in (s01, s02, s03, s04, s05, s06, s07, s08, s09, s10,
               s11, s12, s13, s14, s15, s16, s17, s18, s19, s20):
        fn()
    prs.save(OUT)
    print("OK ->", OUT)
