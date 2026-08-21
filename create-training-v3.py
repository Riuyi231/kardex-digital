from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blue = RGBColor(0x1a, 0x73, 0xe8)
dark = RGBColor(0x1a, 0x1c, 0x23)
white = RGBColor(0xff, 0xff, 0xff)
gray = RGBColor(0x9a, 0x9a, 0x9a)
light = RGBColor(0xf0, 0xf4, 0xf9)
green = RGBColor(0x0d, 0x90, 0x4f)
red = RGBColor(0xea, 0x43, 0x35)
orange = RGBColor(0xf9, 0xab, 0x00)
soft_blue = RGBColor(0xd3, 0xe3, 0xfd)

SCREENSHOTS = r"C:\Users\STIVEN\Documents\Default Project\nexalert-screenshots"
LOGO = r"C:\Users\STIVEN\Documents\Default Project\reportes-equipos\build\icon-512.png"

def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def txt(slide, l, t, w, h, text, size=18, color=dark, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf

def multi_txt(slide, l, t, w, h, lines, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, (text, size, color, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(6)
        p.alignment = align
    return tf

def bar(slide, color=blue):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()

def bottom_bar(slide, color=blue):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.38), Inches(13.333), Inches(0.12))
    shape.fill.solid(); shape.fill.fore_color.rgb = color; shape.line.fill.background()

def add_screenshot(slide, filename, left, top, width, height, border=True):
    path = os.path.join(SCREENSHOTS, filename)
    if not os.path.exists(path):
        txt(slide, left, top, width, height, f"[{filename}]", 12, gray, False, PP_ALIGN.CENTER)
        return
    if border:
        shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left-0.03), Inches(top-0.03), Inches(width+0.06), Inches(height+0.06))
        shadow.fill.solid(); shadow.fill.fore_color.rgb = RGBColor(0xd0, 0xd0, 0xd0); shadow.line.fill.background()
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width), Inches(height))

def monitor_frame(slide, left, top, width, height, img_path=None):
    stand = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + width/2 - 0.2), Inches(top + height), Inches(0.4), Inches(0.2))
    stand.fill.solid(); stand.fill.fore_color.rgb = RGBColor(0x55, 0x55, 0x55); stand.line.fill.background()
    base = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + width/2 - 0.4), Inches(top + height + 0.2), Inches(0.8), Inches(0.06))
    base.fill.solid(); base.fill.fore_color.rgb = RGBColor(0x55, 0x55, 0x55); base.line.fill.background()
    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    frame.fill.solid(); frame.fill.fore_color.rgb = RGBColor(0x2d, 0x2d, 0x2d); frame.line.fill.background()
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(left+0.05), Inches(top+0.05), Inches(width-0.1), Inches(height-0.1))

def phone_frame(slide, left, top, width, height):
    body = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    body.fill.solid(); body.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x1a); body.line.fill.background()
    notch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left + width*0.3), Inches(top), Inches(width*0.4), Inches(0.15))
    notch.fill.solid(); notch.fill.fore_color.rgb = RGBColor(0x1a, 0x1a, 0x1a); notch.line.fill.background()
    inner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left+0.06), Inches(top+0.12), Inches(width-0.12), Inches(height-0.18))
    inner.fill.solid(); inner.fill.fore_color.rgb = white; inner.line.fill.background()
    return inner

def step_card(slide, left, top, num, title, desc, color):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(3.2), Inches(2.8))
    card.fill.solid(); card.fill.fore_color.rgb = white
    card.line.color.rgb = color; card.line.width = Pt(2)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 1.1), Inches(top + 0.2), Inches(1), Inches(1))
    circ.fill.solid(); circ.fill.fore_color.rgb = color; circ.line.fill.background()
    txt(slide, left + 1.1, top + 0.25, 1, 0.8, str(num), 28, white, True, PP_ALIGN.CENTER)
    txt(slide, left + 0.2, top + 1.3, 2.8, 0.5, title, 16, color, True, PP_ALIGN.CENTER)
    txt(slide, left + 0.2, top + 1.8, 2.8, 0.8, desc, 12, gray, False, PP_ALIGN.CENTER)

def bullets(slide, l, t, w, h, items, size=16, color=dark, spacing=8):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(spacing)
    return tf

# ===== SLIDE 1: PORTADA =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, dark)
bar(s); bottom_bar(s)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5))
txt(s, 1, 3, 11.3, 1, "CAPACITACION NEXALERT", 48, white, True, PP_ALIGN.CENTER)
txt(s, 1, 4.2, 11.3, 0.8, "Guia completa para gerentes y tecnicos", 22, gray, False, PP_ALIGN.CENTER)
txt(s, 1, 5.8, 11.3, 0.6, "NEXUS  |  2026", 14, gray, False, PP_ALIGN.CENTER)

# ===== SLIDE 2: QUE ES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Que es NexAlert?", 36, dark, True)
txt(s, 0.8, 1.2, 12, 0.8, "Sistema de gestion de reportes de fallas de equipos", 18, gray)
monitor_frame(s, 0.5, 2.2, 6, 4, os.path.join(SCREENSHOTS, "01-panel.png"))
txt(s, 7, 2.2, 5.8, 0.5, "Una sola plataforma conecta:", 18, dark, True)
items = [
    "\U0001F4BB  App de escritorio para la oficina",
    "\U0001F4F1  App movil para tecnicos en campo",
    "\u2601\ufe0f  Servidor en la nube (Google Cloud)",
    "\U0001F4AC  Integracion con WhatsApp",
    "\U0001F514  Notificaciones push en tiempo real",
    "\U0001F504  Sincronizacion automatica",
]
for i, item in enumerate(items):
    txt(s, 7.2, 2.9 + i*0.5, 5.5, 0.45, item, 14, dark)

# ===== SLIDE 3: ARQUITECTURA =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, light)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Como funciona el sistema", 36, dark, True)

# Desktop
monitor_frame(s, 0.3, 1.8, 4, 3, os.path.join(SCREENSHOTS, "01-panel.png"))
txt(s, 0.3, 5, 4, 0.4, "\U0001F4BB  Desktop - Gerente", 14, blue, True, PP_ALIGN.CENTER)
txt(s, 0.3, 5.4, 4, 0.5, "Panel, clientes, reportes,\ntecnicos, mensajes, ajustes", 11, gray, False, PP_ALIGN.CENTER)

# Arrow
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(3.2), Inches(1.2), Inches(0.5))
arrow.fill.solid(); arrow.fill.fore_color.rgb = blue; arrow.line.fill.background()

# Server
srv = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.9), Inches(1.8), Inches(2.5), Inches(3))
srv.fill.solid(); srv.fill.fore_color.rgb = white; srv.line.color.rgb = blue; srv.line.width = Pt(2)
txt(s, 5.9, 2, 2.5, 0.5, "\u2601\ufe0f  Servidor", 16, blue, True, PP_ALIGN.CENTER)
txt(s, 6.1, 2.6, 2.1, 2, "API REST\nSQLite\nFirebase Push\nAuto-updates\nSync bidireccional\nGCP us-central1", 11, dark, False, PP_ALIGN.CENTER)

# Arrow
arrow2 = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.6), Inches(3.2), Inches(1.2), Inches(0.5))
arrow2.fill.solid(); arrow2.fill.fore_color.rgb = blue; arrow2.line.fill.background()

# Mobile
phone_inner = phone_frame(s, 10, 1.8, 2.5, 3)
txt(s, 10.2, 2.1, 2.1, 0.4, "\U0001F4F1  Movil", 14, blue, True, PP_ALIGN.CENTER)
txt(s, 10.2, 2.6, 2.1, 2, "Reportes\nFoto + GPS\nBiometria\nOffline\nWhatsApp\nNotificaciones", 11, dark, False, PP_ALIGN.CENTER)
txt(s, 10, 5, 2.5, 0.4, "Android (Capacitor)", 11, blue, True, PP_ALIGN.CENTER)

# Bidirectional arrows
arr_l = s.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(4.5), Inches(4), Inches(1.2), Inches(0.5))
arr_l.fill.solid(); arr_l.fill.fore_color.rgb = green; arr_l.line.fill.background()

# ===== SLIDE 4: ROLES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Quien usa que?", 36, dark, True)

# Gerente card
g_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2))
g_card.fill.solid(); g_card.fill.fore_color.rgb = soft_blue; g_card.line.fill.background()
txt(s, 0.8, 1.7, 5.2, 0.5, "\U0001F454  Gerente", 22, blue, True)
txt(s, 0.8, 2.3, 5.2, 0.4, "Desde la PC (escritorio):", 14, dark, True)
items_g = [
    "Ver panel de estadisticas",
    "Gestionar clientes y equipos",
    "Crear y asignar reportes",
    "Exportar a PDF y Excel",
    "Enviar por WhatsApp",
]
for i, b in enumerate(items_g):
    txt(s, 1.1, 2.8 + i*0.4, 4.8, 0.35, "\u2714  " + b, 13, dark)

txt(s, 0.8, 4.8, 5.2, 0.4, "Desde el celular:", 14, dark, True)
items_gm = [
    "Monitorear tecnicos en tiempo real",
    "Ver reportes asignados",
    "Cambiar estados de reportes",
    "Desactivar tecnicos",
]
for i, b in enumerate(items_gm):
    txt(s, 1.1, 5.3 + i*0.4, 4.8, 0.35, "\u2714  " + b, 13, dark)

# Tecnico card
t_card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(1.5), Inches(5.8), Inches(5.2))
t_card.fill.solid(); t_card.fill.fore_color.rgb = RGBColor(0xe8, 0xf5, 0xe9); t_card.line.fill.background()
txt(s, 7.3, 1.7, 5.2, 0.5, "\U0001F477  Tecnico", 22, green, True)
txt(s, 7.3, 2.3, 5.2, 0.4, "Desde el celular (app movil):", 14, dark, True)
items_t = [
    "Ver reportes asignados",
    "Crear reportes con foto y GPS",
    "Cambiar estado de reportes",
    "Compartir por WhatsApp",
    "Trabajar offline",
    "Biometria para entrar rapido",
    "Recibir notificaciones push",
]
for i, b in enumerate(items_t):
    txt(s, 7.6, 2.8 + i*0.4, 4.8, 0.35, "\u2714  " + b, 13, dark)

# ===== SLIDE 5: DESKTOP - PANEL =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Desktop: Panel Principal", 36, dark, True)
add_screenshot(s, "01-panel.png", 0.3, 1.3, 8, 5.3)
txt(s, 8.8, 1.5, 4.2, 0.5, "Lo que ves aqui:", 18, blue, True)
bullets(s, 8.8, 2.2, 4.2, 4, [
    "Estadisticas en tiempo real",
    "Total de reportes activos",
    "Reportes por estado",
    "Reportes por tecnico",
    "Acceso rapido a todas las funciones",
    "Barra de actualizaciones",
    "Sincronizacion con servidor",
], 13, dark, 10)

# ===== SLIDE 6: DESKTOP - REPORTES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Desktop: Gestion de Reportes", 36, dark, True)
add_screenshot(s, "02-reportes.png", 0.3, 1.3, 8, 5.3)
txt(s, 8.8, 1.5, 4.2, 0.5, "Funciones:", 18, blue, True)
bullets(s, 8.8, 2.2, 4.2, 4, [
    "Ver todos los reportes",
    "Filtrar por estado, cliente, fecha",
    "Buscar por texto",
    "Crear nuevo reporte",
    "Cambiar estado con un click",
    "Asignar tecnico",
    "Exportar a PDF y Excel",
    "Enviar por WhatsApp",
], 13, dark, 10)

# ===== SLIDE 7: DESKTOP - CLIENTES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Desktop: Clientes y Equipos", 36, dark, True)
add_screenshot(s, "04-clientes.png", 0.3, 1.3, 8, 5.3)
txt(s, 8.8, 1.5, 4.2, 0.5, "Gestion de clientes:", 18, blue, True)
bullets(s, 8.8, 2.2, 4.2, 2.5, [
    "Agregar nuevos clientes",
    "Editar datos de contacto",
    "Asociar grupos de WhatsApp",
    "Ver historial de reportes",
], 13, dark, 10)
txt(s, 8.8, 4.2, 4.2, 0.5, "Gestion de equipos:", 18, blue, True)
bullets(s, 8.8, 4.8, 4.2, 2, [
    "Registrar equipos por cliente",
    "Marca, modelo, serial",
    "Historial de fallas",
], 13, dark, 10)

# ===== SLIDE 8: DESKTOP - TECNICOS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Desktop: Tecnicos", 36, dark, True)
add_screenshot(s, "05-tecnicos.png", 0.3, 1.3, 8, 5.3)
txt(s, 8.8, 1.5, 4.2, 0.5, "Gestion de tecnicos:", 18, blue, True)
bullets(s, 8.8, 2.2, 4.2, 4, [
    "Agregar tecnicos nuevos",
    "Editar nombre y telefono",
    "Asignar rol (tecnico/gerente)",
    "Importar desde archivo",
    "Activar / desactivar",
    "Cada tecnico tiene su usuario",
], 13, dark, 10)

# ===== SLIDE 9: DESKTOP - MENSAJES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Desktop: Mensajes WhatsApp", 36, dark, True)
add_screenshot(s, "07-mensajes.png", 0.3, 1.3, 8, 5.3)
txt(s, 8.8, 1.5, 4.2, 0.5, "Integracion WhatsApp:", 18, blue, True)
bullets(s, 8.8, 2.2, 4.2, 4, [
    "Chatear con clientes directamente",
    "Ver grupos de WhatsApp",
    "Enviar reportes formateados",
    "Recibir mensajes de clientes",
    "Crear reportes desde chats",
    "Busqueda de conversaciones",
], 13, dark, 10)

# ===== SLIDE 10: DESKTOP - NUEVO REPORTE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Desktop: Crear Reporte", 36, dark, True)
add_screenshot(s, "03-nuevo-reporte.png", 0.3, 1.3, 8, 5.3)
txt(s, 8.8, 1.5, 4.2, 0.5, "Pasos:", 18, blue, True)
bullets(s, 8.8, 2.2, 4.2, 4, [
    "1. Selecciona el cliente",
    "2. Selecciona el equipo",
    "3. Describe el problema",
    "4. Asigna prioridad",
    "5. Asigna tecnico (opcional)",
    "6. Guarda el reporte",
    "7. Se envia por WhatsApp",
], 13, dark, 10)

# ===== SLIDE 11: MOBILE - LOGIN =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, light)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "App Movil: Primer Inicio", 36, dark, True)

phone_inner = phone_frame(s, 1, 1.5, 3.5, 5.5)
txt(s, 1.3, 1.8, 2.9, 0.4, "NexAlert", 18, blue, True, PP_ALIGN.CENTER)
# Login mockup
for i, (label, placeholder) in enumerate([("Usuario", "Tu usuario"), ("Contrasena", "Tu contrasena")]):
    y = 2.6 + i * 0.9
    txt(s, 1.4, y, 2.7, 0.3, label, 10, gray)
    box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(y+0.3), Inches(2.7), Inches(0.35))
    box.fill.solid(); box.fill.fore_color.rgb = light; box.line.fill.background()

btn = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(4.7), Inches(2.7), Inches(0.4))
btn.fill.solid(); btn.fill.fore_color.rgb = blue; btn.line.fill.background()
txt(s, 1.4, 4.72, 2.7, 0.35, "Entrar", 12, white, True, PP_ALIGN.CENTER)

bio = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.4), Inches(5.3), Inches(2.7), Inches(0.4))
bio.fill.solid(); bio.fill.fore_color.rgb = white; bio.line.color.rgb = blue; bio.line.width = Pt(1)
txt(s, 1.4, 5.32, 2.7, 0.35, "Entrar con biometria", 11, blue, True, PP_ALIGN.CENTER)

txt(s, 1.4, 5.9, 2.7, 0.3, "\U0001F511 Huella o rostro", 9, gray, False, PP_ALIGN.CENTER)

# Steps
txt(s, 5.5, 1.5, 7, 0.5, "Como entrar:", 22, dark, True)
step_card(s, 5.5, 2.2, "1", "Escribe usuario", "Tu gerente te da\nel usuario y password", blue)
step_card(s, 9.2, 2.2, "2", "Guarda biometria", "La primera vez guarda\nhuella para despues", green)
step_card(s, 5.5, 5.2, "3", "Entra rapido", "Toca tu dedo\ny listo!", green)
step_card(s, 9.2, 5.2, "4", "Listo!", "Ya puedes ver\ntus reportes", orange)

# ===== SLIDE 12: MOBILE - REPORTES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "App Movil: Mis Reportes", 36, dark, True)

phone_inner = phone_frame(s, 0.5, 1.3, 3.8, 5.8)
txt(s, 0.8, 1.5, 3.2, 0.3, "NexAlert", 12, blue, True, PP_ALIGN.CENTER)
txt(s, 0.8, 1.9, 3.2, 0.3, "Mis reportes", 10, dark, True, PP_ALIGN.CENTER)
# Search bar
sb = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.3), Inches(3.2), Inches(0.3))
sb.fill.solid(); sb.fill.fore_color.rgb = light; sb.line.fill.background()
txt(s, 0.9, 2.32, 3, 0.25, "Buscar cliente, equipo...", 8, gray)
# Report cards
for i, (cliente, equipo, estado, color_e) in enumerate([
    ("tacos El Güero", "Freidora #3", "Abierto", red),
    ("Cafe Azul", "Computo #1", "En proceso", orange),
    ("Farmacia Lopez", "Refrigerador", "Resuelto", green),
]):
    y = 2.8 + i * 1.2
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(y), Inches(3.2), Inches(1))
    card.fill.solid(); card.fill.fore_color.rgb = white; card.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0); card.line.width = Pt(0.5)
    dot = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.9), Inches(y+0.15), Inches(0.15), Inches(0.15))
    dot.fill.solid(); dot.fill.fore_color.rgb = color_e; dot.line.fill.background()
    txt(s, 1.15, y+0.05, 2, 0.25, cliente, 8, dark, True)
    txt(s, 1.15, y+0.3, 2, 0.2, equipo, 7, gray)
    badge = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.2), Inches(y+0.1), Inches(0.7), Inches(0.25))
    badge.fill.solid(); badge.fill.fore_color.rgb = color_e; badge.line.fill.background()
    txt(s, 3.2, y+0.1, 0.7, 0.25, estado.split()[0], 6, white, True, PP_ALIGN.CENTER)

txt(s, 5, 1.3, 7.8, 0.5, "Pantalla principal del tecnico:", 20, dark, True)
bullets(s, 5, 2, 7.8, 4.5, [
    "Ve solo sus reportes asignados",
    "Busca por cliente, equipo o problema",
    "Filtra por estado (abierto, proceso, resuelto...)",
    "Filtra por prioridad (urgente, alta, normal, baja)",
    "Filtra por fechas",
    "Toca un reporte para ver detalles",
    "El boton + crea un reporte nuevo",
], 14, dark, 12)

# ===== SLIDE 13: MOBILE - CREAR REPORTE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "App Movil: Crear Reporte", 36, dark, True)

# Steps with visual
for i, (num, title, desc, col) in enumerate([
    ("1", "Selecciona", "Cliente y equipo", blue),
    ("2", "Describe", "El problema", orange),
    ("3", "Toma foto", "Del equipo danado", green),
    ("4", "Guarda", "Se envia solo", green),
]):
    left = 0.5 + i * 3.2
    step_card(s, left, 1.5, num, title, desc, col)

txt(s, 0.8, 4.7, 12, 0.5, "Caracteristicas del reporte movil:", 18, dark, True)
bullets(s, 0.8, 5.3, 12, 2, [
    "La ubicacion GPS se captura automaticamente  |  Puedes tomar varias fotos  |  Asignar prioridad (baja, normal, alta, urgente)",
    "Si no hay internet, se guarda local y se sincroniza despues  |  El tecnico asignado recibe notificacion push",
], 13, dark, 10)

# ===== SLIDE 14: ESTADOS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Estados de un reporte", 36, dark, True)
txt(s, 0.8, 1.1, 12, 0.5, "Todo reporte pasa por estos estados:", 16, gray)

estados = [
    ("Abierto", "Reporte nuevo,\nsin atender", red, "\u274C"),
    ("En Proceso", "El tecnico\nya esta trabajando", orange, "\u23F3"),
    ("Espera Repuesto", "Necesita una\npieza para reparar", RGBColor(0xf9,0xab,0x00), "\u23F3"),
    ("Espera Cliente", "Se necesita\ninfo del cliente", gray, "\u23F3"),
    ("Resuelto", "Equipo\nreparado", green, "\u2705"),
]
for i, (nome, desc, cor, emoji) in enumerate(estados):
    left = 0.3 + i * 2.55
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(1.8), Inches(2.3), Inches(3))
    card.fill.solid(); card.fill.fore_color.rgb = light; card.line.fill.background()
    top_bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(1.8), Inches(2.3), Inches(0.06))
    top_bar.fill.solid(); top_bar.fill.fore_color.rgb = cor; top_bar.line.fill.background()
    txt(s, left, 2.1, 2.3, 0.6, emoji, 28, cor, True, PP_ALIGN.CENTER)
    txt(s, left + 0.1, 2.7, 2.1, 0.5, nome, 15, cor, True, PP_ALIGN.CENTER)
    txt(s, left + 0.1, 3.3, 2.1, 1, desc, 12, gray, False, PP_ALIGN.CENTER)

# Flow arrows
for i in range(4):
    x = 2.4 + i * 2.55
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(3), Inches(0.3), Inches(0.3))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = RGBColor(0xcc, 0xcc, 0xcc); arrow.line.fill.background()

txt(s, 0.8, 5.2, 12, 0.5, "Despues de 24 horas en 'Resuelto', el reporte se archiva automaticamente", 14, gray, False, PP_ALIGN.CENTER)
txt(s, 0.8, 5.7, 12, 0.5, "Los reportes archivados se pueden restaurar si es necesario", 14, blue, False, PP_ALIGN.CENTER)

# ===== SLIDE 15: WHATSAPP =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Integracion WhatsApp", 36, dark, True)

phone_inner = phone_frame(s, 0.5, 1.3, 3.5, 5.5)
txt(s, 0.8, 1.5, 2.9, 0.3, "WhatsApp", 12, green, True, PP_ALIGN.CENTER)
# Chat bubble
bubble = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2), Inches(2.9), Inches(3.5))
bubble.fill.solid(); bubble.fill.fore_color.rgb = RGBColor(0xdc, 0xf8, 0xc6); bubble.line.fill.background()
txt(s, 0.9, 2.1, 2.7, 3.3, "REPORTE DE FALLA\n\nEmpresa: tacos El Güero\nEquipo: Freidora #3\nFecha: 2026-08-20\nProblema: No enciende,\nse escucha ruido raro\n\n--- NexAlert", 9, dark)

txt(s, 5, 1.5, 7.8, 0.5, "Como funciona:", 22, dark, True)
bullets(s, 5, 2.2, 7.8, 4.5, [
    "El gerente vincula grupos de WhatsApp a cada cliente",
    "Cuando se crea un reporte, se envia automatico al grupo",
    "El mensaje incluye: empresa, equipo, fecha y problema",
    "Desde el movil: toca 'Compartir' en el reporte",
    "Desde el desktop: selecciona el grupo y dale enviar",
    "El cliente recibe el reporte al instante",
    "Tambien puedes chatear directamente con el cliente",
], 14, dark, 12)

# ===== SLIDE 16: OFFLINE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, light)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Modo Offline", 36, dark, True)
txt(s, 0.8, 1.1, 12, 0.5, "La app funciona sin conexion a internet", 18, gray)

for i, (num, title, desc, col) in enumerate([
    ("1", "Sin internet", "Crea tu reporte\nnormalmente", red),
    ("2", "Guardado local", "Se guarda en\ntu celular", orange),
    ("3", "Sincronizacion", "Cuando vuelva\nel internet, se manda", green),
]):
    left = 0.8 + i * 4.2
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(2), Inches(3.5), Inches(3.5))
    card.fill.solid(); card.fill.fore_color.rgb = white; card.line.fill.background()
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 1.2), Inches(2.3), Inches(1.1), Inches(1.1))
    circ.fill.solid(); circ.fill.fore_color.rgb = col; circ.line.fill.background()
    txt(s, left + 1.2, 2.4, 1.1, 0.9, num, 32, white, True, PP_ALIGN.CENTER)
    txt(s, left + 0.3, 3.6, 2.9, 0.5, title, 18, col, True, PP_ALIGN.CENTER)
    txt(s, left + 0.3, 4.2, 2.9, 1, desc, 14, gray, False, PP_ALIGN.CENTER)
    if i < 2:
        arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(left + 3.7), Inches(3.3), Inches(0.4), Inches(0.4))
        arrow.fill.solid(); arrow.fill.fore_color.rgb = col; arrow.line.fill.background()

txt(s, 0.8, 6, 12, 0.5, "Nunca pierdes un reporte aunque no tengas signal", 16, green, True, PP_ALIGN.CENTER)

# ===== SLIDE 17: NOTIFICACIONES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s); bottom_bar(s)
txt(s, 0.8, 0.4, 11, 0.7, "Notificaciones Push", 36, dark, True)

phone_inner = phone_frame(s, 0.5, 1.5, 3.5, 5.5)
# Notification mockup
notif = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.8), Inches(3.1), Inches(1.2))
notif.fill.solid(); notif.fill.fore_color.rgb = white; notif.line.color.rgb = RGBColor(0xe0, 0xe0, 0xe0); notif.line.width = Pt(1)
txt(s, 0.9, 1.9, 2.7, 0.25, "NexAlert", 9, blue, True)
txt(s, 0.9, 2.2, 2.7, 0.25, "Reporte asignado", 10, dark, True)
txt(s, 0.9, 2.5, 2.7, 0.3, "tacos El Güero - Freidora #3\nNo enciende, se escucha ruido", 8, gray)

txt(s, 5, 1.5, 7.8, 0.5, "Cuando se notifica:", 22, dark, True)
bullets(s, 5, 2.2, 7.8, 4, [
    "Cuando un reporte se asigna a un tecnico",
    "Cuando hay un cambio de estado",
    "Cuando el gerente reasigna un reporte",
    "El tecnico recibe: titulo + descripcion",
    "Al tocar la notificacion, se abre el reporte",
    "Funciona aunque la app este cerrada",
], 14, dark, 12)

txt(s, 5, 5.5, 7.8, 0.5, "Configuracion:", 18, dark, True)
txt(s, 5, 6, 7.8, 0.5, "Ajustes del celular > NexAlert > Notificaciones > Activar", 13, gray)

# ===== SLIDE 18: TIPS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, light)
bar(s, green); bottom_bar(s, green)
txt(s, 0.8, 0.4, 11, 0.7, "Tips para el dia a dia", 36, dark, True)

tips = [
    ("\U0001F4F7", "Toma foto del equipo danado"),
    ("\U0001F4DD", "Describe el problema con detalle"),
    ("\U0001F504", "Cambia el estado cuando avances"),
    ("\U0001F511", "Usa biometria para entrar rapido"),
    ("\U0001F4F6", "Sin internet? No te preocupes"),
    ("\U0001F514", "Revisa las notificaciones push"),
    ("\U0001F4E4", "Comparte el reporte al cliente"),
    ("\u2753", "Dudas? Preguntale al gerente"),
]
for i, (emoji, tip) in enumerate(tips):
    col = i % 2
    row = i // 2
    left = 0.8 + col * 6.2
    top = 1.4 + row * 1.4
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(5.8), Inches(1.15))
    card.fill.solid(); card.fill.fore_color.rgb = white; card.line.fill.background()
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left + 0.2), Inches(top + 0.2), Inches(0.7), Inches(0.7))
    circ.fill.solid(); circ.fill.fore_color.rgb = green; circ.line.fill.background()
    txt(s, left + 0.2, top + 0.22, 0.7, 0.6, emoji, 20, white, True, PP_ALIGN.CENTER)
    txt(s, left + 1.1, top + 0.3, 4.4, 0.5, tip, 15, dark)

# ===== SLIDE 19: PROBLEMAS =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, white)
bar(s, orange); bottom_bar(s, orange)
txt(s, 0.8, 0.4, 11, 0.7, "Problemas comunes y soluciones", 36, dark, True)

problemas = [
    ("No puedo entrar", "Verifica tu usuario y contrasena. Si falla 3 veces, reinicia la app.", red),
    ("No me sale la biometria", "Ve a Configuracion y activala. Asegurate de tener huella registrada en el celular.", orange),
    ("No se envia el reporte", "Revisa tu conexion. El reporte se guardo local y se enviara cuando vuelva el internet.", orange),
    ("No recibo notificaciones", "Ve a Ajustes > NexAlert > Notificaciones y activalas.", red),
    ("La app esta lenta", "Cierra otras apps. Si persiste, reinicia el celular.", gray),
]
for i, (prob, sol, col) in enumerate(problemas):
    top = 1.3 + i * 1.2
    pshape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(top), Inches(5), Inches(1))
    pshape.fill.solid()
    pshape.fill.fore_color.rgb = RGBColor(0xfe, 0xe0, 0xe0) if col == red else (RGBColor(0xff, 0xf3, 0xe0) if col == orange else light)
    pshape.line.fill.background()
    txt(s, 0.8, top + 0.1, 4.5, 0.3, prob, 13, col, True)
    txt(s, 0.8, top + 0.45, 4.5, 0.45, sol, 10, dark)
    arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(5.7), Inches(top + 0.3), Inches(0.5), Inches(0.35))
    arrow.fill.solid(); arrow.fill.fore_color.rgb = green; arrow.line.fill.background()
    sshape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(top), Inches(6.4), Inches(1))
    sshape.fill.solid(); sshape.fill.fore_color.rgb = RGBColor(0xe6, 0xf4, 0xe9); sshape.line.fill.background()
    txt(s, 6.7, top + 0.3, 5.8, 0.4, "Solucion", 11, green, True)

# ===== SLIDE 20: CIERRE =====
s = prs.slides.add_slide(prs.slide_layouts[6])
bg(s, dark)
bar(s); bottom_bar(s)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(5.9), Inches(1.2), Inches(1.5), Inches(1.5))
txt(s, 1, 3, 11.3, 1, "Listo para empezar!", 48, white, True, PP_ALIGN.CENTER)
txt(s, 1, 4.2, 11.3, 0.8, "Descarga la app, inicia sesion y crea tu primer reporte", 20, gray, False, PP_ALIGN.CENTER)
txt(s, 1, 5.5, 11.3, 0.6, "Dudas? Consulta al gerente", 14, gray, False, PP_ALIGN.CENTER)

output = os.path.join(os.path.expanduser("~"), "Desktop", "NexAlert-Capacitacion.pptx")
prs.save(output)
print(f"Capacitacion creada: {output}")
